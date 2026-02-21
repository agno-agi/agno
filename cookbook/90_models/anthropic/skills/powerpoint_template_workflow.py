"""
Agno Workflow: PowerPoint Template Generation Pipeline.

A sequential workflow that generates presentations using Claude's pptx skill,
intelligently adds AI-generated images via NanoBanana (powered by Gemini), and
applies a custom .pptx template for professional styling.

Prerequisites:
- uv pip install agno anthropic python-pptx google-genai pillow
- export ANTHROPIC_API_KEY="your_api_key_here"
- export GOOGLE_API_KEY="your_google_api_key_here"
- A .pptx template file

Usage:
    # Basic usage with a template:
    .venvs/demo/bin/python cookbook/90_models/anthropic/skills/powerpoint_template_workflow.py \\
        --template my_template.pptx

    # Full options: custom prompt, output path, verbose logging:
    .venvs/demo/bin/python cookbook/90_models/anthropic/skills/powerpoint_template_workflow.py \\
        -t my_template.pptx -o report.pptx -p "Create a 5-slide AI trends presentation" -v

    # Disable streaming (more reliable for shorter prompts, may timeout on complex ones):
    .venvs/demo/bin/python cookbook/90_models/anthropic/skills/powerpoint_template_workflow.py \\
        -t my_template.pptx --no-stream

    # Skip AI image generation entirely:
    .venvs/demo/bin/python cookbook/90_models/anthropic/skills/powerpoint_template_workflow.py \\
        -t my_template.pptx --no-images

CLI Flags:
    --template, -t   Path to the .pptx template file (required).
    --output, -o     Output filename (default: presentation_from_template.pptx).
    --prompt, -p     Custom prompt for the presentation content.
    --no-images      Skip AI image generation (Steps 2 and 3).
    --no-stream      Disable streaming mode for Claude agent (more reliable for
                     shorter prompts, but may timeout on complex presentations).
    --verbose, -v    Enable verbose/debug logging for troubleshooting.
"""

import argparse
import copy
import json
import os
import shutil
import sys
import traceback
from dataclasses import dataclass, field
from io import BytesIO
from typing import Dict, List

from agno.agent import Agent
from agno.models.anthropic import Claude
from agno.models.google import Gemini
from agno.run.agent import RunOutput
from agno.tools.nano_banana import NanoBananaTools
from agno.workflow.step import Step
from agno.workflow.types import StepInput, StepOutput
from agno.workflow.workflow import Workflow
from anthropic import Anthropic
from file_download_helper import download_skill_files
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

# Module-level verbose flag (set from CLI args)
VERBOSE = False


# ---------------------------------------------------------------------------
# Pydantic Models for Structured Data Flow
# ---------------------------------------------------------------------------


class SlideImageDecision(BaseModel):
    """Decision about whether a slide needs an AI-generated image."""

    slide_index: int = Field(description="Zero-based index of the slide")
    needs_image: bool = Field(
        description="Whether this slide would benefit from an image"
    )
    image_prompt: str = Field(
        default="",
        description="If needs_image is True, a detailed prompt for generating the image. "
        "Should describe a professional, clean illustration suitable for a business presentation. "
        "If needs_image is False, leave empty.",
    )
    reasoning: str = Field(
        default="",
        description="Brief explanation of why the slide does or does not need an image",
    )


class ImagePlan(BaseModel):
    """Plan for which slides need AI-generated images."""

    decisions: List[SlideImageDecision] = Field(
        description="List of image decisions, one per slide"
    )


# ---------------------------------------------------------------------------
# Dataclasses for Slide Content (reused from agent_with_powerpoint_template)
# ---------------------------------------------------------------------------


@dataclass
class TableData:
    """Extracted table data with position."""

    rows: list
    left: int
    top: int
    width: int
    height: int


@dataclass
class ImageData:
    """Extracted image data with position."""

    blob: bytes
    left: int
    top: int
    width: int
    height: int
    content_type: str = "image/png"


@dataclass
class ChartExtract:
    """Extracted chart data with position."""

    chart_type: int
    categories: list
    series: list
    left: int
    top: int
    width: int
    height: int


@dataclass
class SlideContent:
    """All extracted content from a single slide."""

    title: str = ""
    subtitle: str = ""
    body_paragraphs: list = field(default_factory=list)
    tables: list = field(default_factory=list)
    images: list = field(default_factory=list)
    charts: list = field(default_factory=list)
    shapes_xml: list = field(default_factory=list)
    # Track picture placeholders detected in the slide layout
    has_image_placeholder: bool = False
    image_placeholder_indices: list = field(default_factory=list)


@dataclass
class ContentArea:
    """Defines the safe content region on a template slide (all values in EMU)."""

    left: int
    top: int
    width: int
    height: int


# ---------------------------------------------------------------------------
# Content Extraction Functions (from agent_with_powerpoint_template.py)
# ---------------------------------------------------------------------------


def _extract_slide_content(slide) -> SlideContent:
    """Extract all content from a slide including text, tables, images, charts, and shapes.

    Detects picture placeholders using multiple strategies for robustness:
      1. int(type) == 18  (the raw OOXML value for PICTURE placeholders)
      2. str(type) contains 'PICTURE (18)'
      3. PP_PLACEHOLDER.PICTURE enum comparison
      4. XML-level fallback: <p:ph type="pic"/>
    """
    content = SlideContent()

    # Detect picture placeholders in the slide layout
    for shape in slide.placeholders:
        ph_fmt = shape.placeholder_format
        if ph_fmt is not None:
            ph_type_val = ph_fmt.type
            # Strategy 1 & 2: enum / int / string comparison
            is_picture_ph = False
            try:
                if ph_type_val is not None and (
                    int(ph_type_val) == 18 or str(ph_type_val) == "PICTURE (18)"
                ):
                    is_picture_ph = True
            except (ValueError, TypeError):
                pass
            # Also try the enum directly
            if not is_picture_ph:
                try:
                    if ph_type_val == PP_PLACEHOLDER.PICTURE:
                        is_picture_ph = True
                except Exception as e:
                    if VERBOSE:
                        print("[VERBOSE] Exception suppressed: %s" % str(e))
            # Strategy 3: XML-level fallback
            if not is_picture_ph:
                nsmap = {
                    "p": "http://schemas.openxmlformats.org/presentationml/2006/main"
                }
                ph_elem = shape._element.find(".//p:ph", nsmap)
                if ph_elem is not None and ph_elem.get("type") == "pic":
                    is_picture_ph = True
            if is_picture_ph:
                content.has_image_placeholder = True
                content.image_placeholder_indices.append(ph_fmt.idx)

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_data)
            content.tables.append(
                TableData(
                    rows=rows_data,
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                )
            )
            continue

        if shape.has_chart:
            try:
                chart = shape.chart
                chart_type_val = chart.chart_type
                categories = []
                series_data = []
                plot = chart.plots[0] if chart.plots else None
                if plot:
                    if plot.categories:
                        categories = list(plot.categories)
                    for series in plot.series:
                        name = series.name if hasattr(series, "name") else ""
                        values = (
                            list(series.values) if hasattr(series, "values") else []
                        )
                        series_data.append((name or "", values))
                if categories or series_data:
                    content.charts.append(
                        ChartExtract(
                            chart_type=chart_type_val,
                            categories=categories,
                            series=series_data,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height,
                        )
                    )
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                ct = shape.image.content_type
                content.images.append(
                    ImageData(
                        blob=blob,
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                        content_type=ct,
                    )
                )
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_xml = copy.deepcopy(shape._element)
            content.shapes_xml.append(group_xml)
            try:
                for grp_shape in shape.shapes:
                    if grp_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            blob = grp_shape.image.blob
                            ct = grp_shape.image.content_type
                            content.images.append(
                                ImageData(
                                    blob=blob,
                                    left=grp_shape.left,
                                    top=grp_shape.top,
                                    width=grp_shape.width,
                                    height=grp_shape.height,
                                    content_type=ct,
                                )
                            )
                        except Exception as e:
                            if VERBOSE:
                                print("[VERBOSE] Exception suppressed: %s" % str(e))
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.has_text_frame:
            text_frame = shape.text_frame
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0:
                    content.title = text_frame.text.strip()
                elif ph_idx == 1:
                    paragraphs = text_frame.paragraphs
                    if len(paragraphs) == 1 and paragraphs[0].level == 0:
                        content.subtitle = paragraphs[0].text.strip()
                    else:
                        for para in paragraphs:
                            text = para.text.strip()
                            if text:
                                content.body_paragraphs.append((text, para.level))
                else:
                    for para in text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            content.body_paragraphs.append((text, para.level))
            else:
                shape_xml = copy.deepcopy(shape._element)
                content.shapes_xml.append(shape_xml)
        elif not shape.is_placeholder:
            shape_xml = copy.deepcopy(shape._element)
            content.shapes_xml.append(shape_xml)

    return content


# ---------------------------------------------------------------------------
# Template Application Functions (from agent_with_powerpoint_template.py)
# ---------------------------------------------------------------------------


def _find_best_layout(template_prs, slide_index: int, total_slides: int):
    """Find the best matching layout from the template for a given slide position."""
    layouts = list(template_prs.slide_layouts)
    if not layouts:
        raise ValueError("Template has no slide layouts")

    layout_names = [(i, layout.name.lower()) for i, layout in enumerate(layouts)]
    is_title_slide = slide_index == 0
    is_last_slide = slide_index == total_slides - 1

    if is_title_slide:
        for i, name in layout_names:
            if "title slide" in name or (
                "title" in name and "content" not in name and "only" not in name
            ):
                return layouts[i]
        for i, name in layout_names:
            if "title" in name:
                return layouts[i]
        return layouts[0]

    if is_last_slide:
        for i, name in layout_names:
            if "blank" in name or "closing" in name or "end" in name:
                return layouts[i]

    for i, name in layout_names:
        if "content" in name or "body" in name or "text" in name:
            return layouts[i]

    for i, name in layout_names:
        if "object" in name or "list" in name:
            return layouts[i]

    if len(layouts) > 1:
        return layouts[1]
    return layouts[0]


def _get_content_area(layout, slide_width: int, slide_height: int) -> ContentArea:
    """Derive the safe content area from a template layout's placeholders.

    Strategy:
    1. Look for a body placeholder (idx=1) -- its position defines the content area.
    2. If no body placeholder, look for any placeholder with idx > 0.
    3. If no placeholders at all, use a default safe margin.

    Args:
        layout: A python-pptx SlideLayout object.
        slide_width: Presentation slide width in EMU.
        slide_height: Presentation slide height in EMU.

    Returns:
        ContentArea with the computed safe region.
    """
    # Try body placeholder first (idx=1)
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == 1:
            return ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )

    # Try any non-title placeholder
    for ph in layout.placeholders:
        if ph.placeholder_format.idx > 0:
            return ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )

    # Default: safe margins (5% left, 25% top, 90% width, 65% height)
    return ContentArea(
        left=int(slide_width * 0.05),
        top=int(slide_height * 0.25),
        width=int(slide_width * 0.90),
        height=int(slide_height * 0.65),
    )


def _fit_to_area(img_width: int, img_height: int, area: ContentArea) -> tuple:
    """Scale dimensions to fit within an area while preserving aspect ratio.

    Args:
        img_width: Original width in EMU.
        img_height: Original height in EMU.
        area: Target ContentArea to fit within.

    Returns:
        Tuple of (left, top, width, height) in EMU, centered in the area.
    """
    if img_width <= 0 or img_height <= 0:
        return area.left, area.top, area.width, area.height

    aspect = img_width / img_height
    area_aspect = area.width / area.height

    if aspect > area_aspect:
        new_width = area.width
        new_height = int(new_width / aspect)
    else:
        new_height = area.height
        new_width = int(new_height * aspect)

    left = area.left + (area.width - new_width) // 2
    top = area.top + (area.height - new_height) // 2

    return left, top, new_width, new_height


def _populate_placeholder_with_format(shape, texts, is_title=False):
    """Populate a placeholder shape with text while preserving template formatting.

    Enables word wrap and attempts to auto-fit text to the placeholder bounds.
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Enable word wrap before anything else
    tf.word_wrap = True

    ref_para = tf.paragraphs[0] if tf.paragraphs else None
    ref_para_xml = None
    ref_run_xml = None

    if ref_para is not None:
        pPr = ref_para._p.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
        )
        if pPr is not None:
            ref_para_xml = copy.deepcopy(pPr)
        if ref_para.runs:
            rPr = ref_para.runs[0]._r.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
            )
            if rPr is not None:
                ref_run_xml = copy.deepcopy(rPr)

    tf.clear()

    if is_title:
        para = tf.paragraphs[0]
        if ref_para_xml is not None:
            para._p.insert(0, copy.deepcopy(ref_para_xml))
        run = para.add_run()
        run.text = texts
        if ref_run_xml is not None:
            run._r.insert(0, copy.deepcopy(ref_run_xml))
    else:
        for i, (text, level) in enumerate(texts):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            if ref_para_xml is not None:
                new_pPr = copy.deepcopy(ref_para_xml)
                if level > 0:
                    new_pPr.set("lvl", str(level))
                para._p.insert(0, new_pPr)
            para.level = level
            run = para.add_run()
            run.text = text
            if ref_run_xml is not None:
                run._r.insert(0, copy.deepcopy(ref_run_xml))

    # Auto-fit text to placeholder bounds
    try:
        max_size = 28 if is_title else 18
        tf.fit_text(font_family="Calibri", max_size=max_size)
    except Exception as e:
        # fit_text requires font metrics; fall back to MSO_AUTO_SIZE
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _transfer_tables(slide, tables, content_area: ContentArea):
    """Transfer extracted table data to a slide, repositioned to the content area."""
    TABLE_CELL_FONT_SIZE = Pt(10)
    TABLE_HEADER_FONT_SIZE = Pt(11)
    num_tables = len(tables)

    for t_idx, td in enumerate(tables):
        num_rows = len(td.rows)
        num_cols = len(td.rows[0]) if td.rows else 0
        if num_rows == 0 or num_cols == 0:
            continue

        # Stack multiple tables vertically within the content area
        if num_tables > 1:
            per_table_height = content_area.height // num_tables
            table_top = content_area.top + (t_idx * per_table_height)
            table_height = per_table_height
        else:
            table_top = content_area.top
            table_height = content_area.height

        table_shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            content_area.left,
            table_top,
            content_area.width,
            table_height,
        )
        table = table_shape.table
        for r_idx, row_data in enumerate(td.rows):
            for c_idx, cell_text in enumerate(row_data):
                if c_idx < num_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_text
                    # Control font size to prevent overflow
                    for para in cell.text_frame.paragraphs:
                        para.font.size = (
                            TABLE_HEADER_FONT_SIZE
                            if r_idx == 0
                            else TABLE_CELL_FONT_SIZE
                        )
                    cell.text_frame.word_wrap = True


def _transfer_images(slide, images, content_area: ContentArea):
    """Transfer extracted images to a slide, scaled to fit the content area."""
    for img in images:
        image_stream = BytesIO(img.blob)
        left, top, width, height = _fit_to_area(img.width, img.height, content_area)
        slide.shapes.add_picture(image_stream, left, top, width, height)


def _transfer_charts(slide, charts, content_area: ContentArea):
    """Transfer extracted chart data to a slide, sized to fill the content area."""
    try:
        from pptx.chart.data import CategoryChartData
    except ImportError:
        return

    num_charts = len(charts)

    for c_idx, cd in enumerate(charts):
        try:
            chart_data = CategoryChartData()
            chart_data.categories = cd.categories
            for series_name, series_values in cd.series:
                clean_values = []
                for v in series_values:
                    if v is None:
                        clean_values.append(0)
                    elif isinstance(v, (int, float)):
                        clean_values.append(v)
                    else:
                        try:
                            clean_values.append(float(v))
                        except (ValueError, TypeError):
                            clean_values.append(0)
                chart_data.add_series(series_name, clean_values)

            # Stack multiple charts vertically within the content area
            if num_charts > 1:
                chart_height = content_area.height // num_charts
                chart_top = content_area.top + (c_idx * chart_height)
            else:
                chart_height = content_area.height
                chart_top = content_area.top

            slide.shapes.add_chart(
                cd.chart_type,
                content_area.left,
                chart_top,
                content_area.width,
                chart_height,
                chart_data,
            )
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))


def _transfer_shapes(slide, shapes_xml):
    """Transfer simple shapes by deep-copying their XML to the target slide."""
    spTree = slide.shapes._spTree
    for shape_elem in shapes_xml:
        cloned = copy.deepcopy(shape_elem)
        existing_ids = [
            int(sp.get("id", 0)) for sp in spTree.iter() if sp.get("id") is not None
        ]
        max_id = max(existing_ids) if existing_ids else 0
        for nv_elem in cloned.iter():
            if nv_elem.tag.endswith("}cNvPr"):
                max_id += 1
                nv_elem.set("id", str(max_id))
        spTree.append(cloned)


def _clear_unused_placeholders(slide, populated_indices: set) -> None:
    """Remove unused placeholder XML elements from slide to prevent ghost text.

    This is the only reliable way to eliminate ALL types of ghost text:
    - "Click to add title"
    - "Click to add text"
    - "Click icon to add picture"
    - Content placeholder insertion icons (table, chart, image, etc.)

    Simply calling tf.clear() is insufficient for picture placeholders and
    content placeholders with embedded icons.  Removing the XML element
    from the shape tree is the nuclear option that works for every case.

    Args:
        slide: The pptx slide object.
        populated_indices: Set of placeholder idx values that were filled with content.
    """
    spTree = slide.shapes._spTree
    elements_to_remove = []

    # Snapshot the collection with list() to avoid proxy/iterator issues
    for shape in list(slide.placeholders):
        ph_idx = shape.placeholder_format.idx
        if ph_idx in populated_indices:
            continue  # This placeholder was populated, keep it

        # Don't remove shapes that have actual content (charts, tables, real images)
        try:
            if hasattr(shape, "has_chart") and shape.has_chart:
                continue
            if hasattr(shape, "has_table") and shape.has_table:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))

        # Mark for removal -- removing from XML is the only reliable cleanup
        elements_to_remove.append(shape._element)

    for element in elements_to_remove:
        spTree.remove(element)


def _populate_slide(
    new_slide,
    content: SlideContent,
    slide_width: int,
    slide_height: int,
    generated_image_bytes: bytes | None = None,
):
    """Transfer all content into a new slide using template-aware positioning.

    Handles text placeholders, tables, charts, images, and shapes. When
    ``generated_image_bytes`` is provided and the slide layout contains a
    picture placeholder, the image is inserted into that placeholder rather
    than being added as a free-floating picture.

    Args:
        new_slide: The new slide created from a template layout.
        content: Extracted SlideContent from the generated slide.
        slide_width: Presentation slide width in EMU.
        slide_height: Presentation slide height in EMU.
        generated_image_bytes: Optional raw image bytes (e.g. from NanoBanana)
            to insert into picture placeholders.
    """
    # Compute content area from the slide's layout
    content_area = _get_content_area(new_slide.slide_layout, slide_width, slide_height)

    # Track populated placeholder indices for cleanup
    populated_indices: set[int] = set()

    title_placed = False
    body_placed = False

    for shape in new_slide.placeholders:
        ph_idx = shape.placeholder_format.idx
        if ph_idx == 0 and content.title:
            _populate_placeholder_with_format(shape, content.title, is_title=True)
            populated_indices.add(ph_idx)
            title_placed = True
        elif ph_idx == 1:
            if content.body_paragraphs:
                _populate_placeholder_with_format(
                    shape, content.body_paragraphs, is_title=False
                )
                populated_indices.add(ph_idx)
                body_placed = True
            elif content.subtitle:
                _populate_placeholder_with_format(
                    shape, content.subtitle, is_title=True
                )
                populated_indices.add(ph_idx)
                body_placed = True

    if not body_placed and content.body_paragraphs:
        for shape in new_slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            if ph_idx > 1 and shape.has_text_frame:
                _populate_placeholder_with_format(
                    shape, content.body_paragraphs, is_title=False
                )
                populated_indices.add(ph_idx)
                body_placed = True
                break

    # Fallback text boxes using content area bounds
    if not title_placed and content.title:
        txBox = new_slide.shapes.add_textbox(
            content_area.left,
            Inches(0.3),
            content_area.width,
            Inches(1.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = content.title
        for para in tf.paragraphs:
            para.font.size = Pt(28)
            para.font.bold = True

    if not body_placed and content.body_paragraphs:
        txBox = new_slide.shapes.add_textbox(
            content_area.left,
            content_area.top,
            content_area.width,
            content_area.height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, (text, level) in enumerate(content.body_paragraphs):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = text
            para.level = level
            para.font.size = Pt(18)
        try:
            tf.fit_text(font_family="Calibri", max_size=18)
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    _transfer_tables(new_slide, content.tables, content_area)
    _transfer_images(new_slide, content.images, content_area)
    _transfer_charts(new_slide, content.charts, content_area)
    _transfer_shapes(new_slide, content.shapes_xml)

    # ------------------------------------------------------------------
    # Insert generated images into picture placeholders.
    # Detection uses both the tracked indices AND XML-level type checking
    # as a fallback (type="pic" or "clipArt" in the XML = picture placeholder).
    # ------------------------------------------------------------------
    if generated_image_bytes:
        # Build a comprehensive set of picture placeholder indices using
        # all three detection strategies: int(18), PP_PLACEHOLDER.PICTURE,
        # and XML-level type="pic" / type="clipArt".
        _pic_ph_indices = set(content.image_placeholder_indices)

        for shape in list(new_slide.placeholders):
            ph_fmt = shape.placeholder_format
            if ph_fmt is not None:
                is_picture_ph = False
                # Strategy 1: int comparison with raw OOXML value 18
                try:
                    if ph_fmt.type is not None and (
                        int(ph_fmt.type) == 18 or str(ph_fmt.type) == "PICTURE (18)"
                    ):
                        is_picture_ph = True
                except (ValueError, TypeError):
                    pass
                # Strategy 2: PP_PLACEHOLDER enum
                if not is_picture_ph:
                    try:
                        if ph_fmt.type == PP_PLACEHOLDER.PICTURE:
                            is_picture_ph = True
                    except Exception as e:
                        if VERBOSE:
                            print("[VERBOSE] Exception suppressed: %s" % str(e))
                # Strategy 3: XML-level fallback
                if not is_picture_ph:
                    try:
                        nsmap = {
                            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                        }
                        ph_elem = shape._element.find(".//p:ph", nsmap)
                        if ph_elem is not None and ph_elem.get("type") in (
                            "pic",
                            "clipArt",
                        ):
                            is_picture_ph = True
                    except Exception as e:
                        if VERBOSE:
                            print("[VERBOSE] Exception suppressed: %s" % str(e))
                if is_picture_ph:
                    _pic_ph_indices.add(ph_fmt.idx)

        # Insert the image into the first available picture placeholder
        for shape in list(new_slide.placeholders):
            ph_idx = shape.placeholder_format.idx
            if ph_idx in _pic_ph_indices and ph_idx not in populated_indices:
                try:
                    image_stream = BytesIO(generated_image_bytes)
                    shape.insert_picture(image_stream)
                    populated_indices.add(ph_idx)
                    break  # Use first picture placeholder
                except Exception as e:
                    if VERBOSE:
                        print("[VERBOSE] Exception suppressed: %s" % str(e))

    # ------------------------------------------------------------------
    # Remove ALL unpopulated placeholder XML elements from the slide shape
    # tree. This is the ONLY reliable way to eliminate ALL types of ghost
    # text ("Click to add title", "Click to add text", "Click icon to add
    # picture", content placeholder icons, etc.).
    # tf.clear() alone is insufficient for picture and content placeholders.
    # ------------------------------------------------------------------
    _clear_unused_placeholders(new_slide, populated_indices)


# ---------------------------------------------------------------------------
# Step 1: Content Generation
# ---------------------------------------------------------------------------


def step_generate_content(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Generate PowerPoint content using Claude's pptx skill.

    Supports both streaming and non-streaming modes (controlled by the
    ``stream`` key in session_state). Streaming is required for long-running
    skill operations (>10 min) but may have issues with provider_data
    propagation. Non-streaming is simpler and more reliable for shorter
    operations but can timeout on complex presentations.

    This step:
    1. Creates a Claude agent with the pptx skill
    2. Runs the user's prompt to generate a presentation
    3. Downloads the generated .pptx file
    4. Extracts SlideContent from each slide
    5. Stores the extracted content in session_state
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)

    user_prompt = step_input.input
    template_path = session_state.get("template_path", "")
    output_dir = session_state.get("output_dir", ".")

    print("=" * 60)
    print("Step 1: Generating presentation content with Claude...")
    print("=" * 60)

    # Build prompt with template context
    enhanced_prompt = user_prompt
    try:
        prs = Presentation(template_path)
        layouts = [layout.name for layout in prs.slide_layouts]
        layout_info = ", ".join(layouts)
        enhanced_prompt = (
            user_prompt + "\n\n"
            "Important structural requirements for template compatibility:\n"
            "- Use one clear title and concise bullet points per slide.\n"
            "- Do not apply custom fonts, colors, or theme styling.\n"
            "- Tables and charts are supported and will be transferred to the template.\n"
            "- Keep tables to max 6 rows x 5 columns.\n"
            "- Use bar, column, line, or pie charts with clearly labeled data.\n"
            "- The template has these available layouts: " + layout_info + ".\n"
            "- Use standard slide ordering: Title Slide, then Content Slides, then Closing."
        )
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))

    # Create the Claude agent
    content_agent = Agent(
        name="Content Generator",
        model=Claude(
            id="claude-sonnet-4-5-20250929",  # Or: "claude-sonnet-4-6", "claude-opus-4-6"
            skills=[{"type": "anthropic", "skill_id": "pptx", "version": "latest"}],
        ),
        instructions=[
            "You are a structured content generator for PowerPoint presentations.",
            "Your output will be extracted and remapped to a corporate template.",
            "",
            "SLIDE STRUCTURE:",
            "- Use exactly one clear, descriptive title per slide.",
            "- Use concise bullet points for body content.",
            "- Limit to 4-6 bullet points per slide, each bullet max ~15 words.",
            "- Keep subtitle text on title slides to a single short line.",
            "- Follow standard slide ordering: Title Slide, then Content Slides, then Closing Slide.",
            "",
            "VISUAL ELEMENTS (use when appropriate):",
            "- Include tables for data comparisons and structured information.",
            "- Keep tables concise: no more than 6 rows and 5 columns.",
            "- Use bar, column, line, or pie charts for data visualization.",
            "",
            "FORMATTING RESTRICTIONS:",
            "- Do NOT apply custom fonts, colors, or theme styling.",
            "- Do NOT use SmartArt or complex nested graphic layouts.",
            "- Do NOT add speaker notes, animations, or transitions.",
        ],
        markdown=True,
    )

    # Determine streaming mode from session_state
    use_stream = session_state.get("stream", True)

    if use_stream:
        # Streaming mode: required for long-running skill operations (>10 min)
        response: RunOutput | None = None
        event_count = 0
        try:
            for event in content_agent.run(
                enhanced_prompt, stream=True, yield_run_output=True
            ):
                event_count += 1
                if isinstance(event, RunOutput):
                    response = event
                    if VERBOSE:
                        print(
                            "[VERBOSE] Received RunOutput after %d events" % event_count
                        )
                elif VERBOSE and event_count <= 5:
                    print(
                        "[VERBOSE] Stream event %d: type=%s"
                        % (event_count, type(event).__name__)
                    )
                elif VERBOSE and event_count == 6:
                    print("[VERBOSE] (Suppressing further stream event logs...)")
        except Exception as e:
            print("[ERROR] Agent streaming failed with exception: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            return StepOutput(
                content="Error: Agent execution failed: %s" % str(e),
                success=False,
                stop=True,
            )
        if VERBOSE:
            print("[VERBOSE] Total stream events received: %d" % event_count)
    else:
        # Non-streaming mode: simpler, more reliable for shorter operations
        response = None
        try:
            response = content_agent.run(enhanced_prompt, stream=False)
        except Exception as e:
            print("[ERROR] Agent execution failed with exception: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            return StepOutput(
                content="Error: Agent execution failed: %s" % str(e),
                success=False,
                stop=True,
            )

    if response is None:
        return StepOutput(
            content="Error: No response received from content agent.",
            success=False,
            stop=True,
        )
    print("\nAgent response received.")

    if VERBOSE:
        # Diagnostic: Inspect the RunOutput
        print("\n[VERBOSE] === Agent Response Diagnostics ===")
        if response is None:
            print("[VERBOSE] Response is None!")
        else:
            print("[VERBOSE] Response type: %s" % type(response).__name__)
            print(
                "[VERBOSE] Number of messages: %d"
                % (len(response.messages) if response.messages else 0)
            )
            if response.messages:
                for m_idx, msg in enumerate(response.messages):
                    print("[VERBOSE] Message %d:" % m_idx)
                    print("[VERBOSE]   Type: %s" % type(msg).__name__)
                    print("[VERBOSE]   Role: %s" % getattr(msg, "role", "N/A"))
                    print(
                        "[VERBOSE]   Has provider_data: %s"
                        % (
                            hasattr(msg, "provider_data")
                            and msg.provider_data is not None
                        )
                    )
                    if hasattr(msg, "provider_data") and msg.provider_data:
                        pd = msg.provider_data
                        print("[VERBOSE]   provider_data type: %s" % type(pd).__name__)
                        if isinstance(pd, dict):
                            print(
                                "[VERBOSE]   provider_data keys: %s" % list(pd.keys())
                            )
                        elif hasattr(pd, "__dict__"):
                            print(
                                "[VERBOSE]   provider_data attrs: %s"
                                % list(vars(pd).keys())[:20]
                            )
                    # Also show content snippet
                    content_text = getattr(msg, "content", None)
                    if content_text and isinstance(content_text, str):
                        print(
                            "[VERBOSE]   Content (first 200 chars): %s"
                            % content_text[:200]
                        )
            # Check for response-level attributes
            for attr_name in [
                "content",
                "tool_calls",
                "extra_data",
                "metrics",
                "model_provider_data",
            ]:
                val = getattr(response, attr_name, "MISSING")
                if val != "MISSING":
                    if isinstance(val, str):
                        print(
                            "[VERBOSE] response.%s (first 200 chars): %s"
                            % (attr_name, val[:200])
                        )
                    elif val is not None:
                        print(
                            "[VERBOSE] response.%s type: %s"
                            % (attr_name, type(val).__name__)
                        )
        print("[VERBOSE] === End Diagnostics ===\n")

    # Download the generated file
    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    generated_file = None

    if response.messages:
        for msg_idx, msg in enumerate(response.messages):
            if VERBOSE:
                print("[VERBOSE] Checking message %d for provider_data..." % msg_idx)
            if hasattr(msg, "provider_data") and msg.provider_data:
                if VERBOSE:
                    print(
                        "[VERBOSE] Message %d has provider_data, calling download_skill_files..."
                        % msg_idx
                    )
                try:
                    files = download_skill_files(
                        msg.provider_data, client, output_dir=output_dir
                    )
                except Exception as e:
                    print(
                        "[ERROR] download_skill_files raised an exception: %s" % str(e)
                    )
                    if VERBOSE:
                        traceback.print_exc()
                    files = None

                if VERBOSE:
                    print("[VERBOSE] download_skill_files returned: %s" % files)

                if files:
                    for f in files:
                        if not f.endswith(".pptx"):
                            if VERBOSE:
                                print("[VERBOSE] Skipping non-.pptx file: %s" % f)
                            continue
                        try:
                            Presentation(f)
                            generated_file = f
                            print("Downloaded: %s" % generated_file)
                            break
                        except Exception as e:
                            print(
                                "[WARNING] File '%s' failed pptx validation: %s"
                                % (f, str(e))
                            )
                            continue
                    if not generated_file:
                        # Fallback: try any file
                        for f in files:
                            try:
                                Presentation(f)
                                generated_file = f
                                print("Downloaded (fallback): %s" % generated_file)
                                break
                            except Exception as e:
                                if VERBOSE:
                                    print(
                                        "[VERBOSE] Fallback file '%s' also failed: %s"
                                        % (f, str(e))
                                    )
                                continue
                elif VERBOSE:
                    print("[VERBOSE] No files returned from download_skill_files")
            elif VERBOSE:
                print("[VERBOSE] Message %d has no provider_data" % msg_idx)
            if generated_file:
                break
    else:
        if VERBOSE:
            print("[VERBOSE] response.messages is empty or None!")

    # Fallback: check response.model_provider_data directly
    # This catches cases where provider_data is on the RunOutput but not on individual messages
    if not generated_file and response.model_provider_data:
        if VERBOSE:
            print(
                "[VERBOSE] Trying fallback: response.model_provider_data = %s"
                % response.model_provider_data
            )
        try:
            files = download_skill_files(
                response.model_provider_data, client, output_dir=output_dir
            )
        except Exception as e:
            print("[ERROR] download_skill_files (fallback) raised: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            files = None
        if files:
            if VERBOSE:
                print("[VERBOSE] Fallback download_skill_files returned: %s" % files)
            for f in files:
                if f.endswith(".pptx"):
                    try:
                        Presentation(f)
                        generated_file = f
                        print(
                            "Downloaded (via model_provider_data fallback): %s"
                            % generated_file
                        )
                        break
                    except Exception as e:
                        print(
                            "[WARNING] Fallback file '%s' failed pptx validation: %s"
                            % (f, str(e))
                        )
            if not generated_file:
                for f in files:
                    try:
                        Presentation(f)
                        generated_file = f
                        print(
                            "Downloaded (via model_provider_data fallback): %s"
                            % generated_file
                        )
                        break
                    except Exception as e:
                        if VERBOSE:
                            print(
                                "[VERBOSE] Fallback file '%s' also failed: %s"
                                % (f, str(e))
                            )

    if not generated_file:
        if VERBOSE:
            print("\n[VERBOSE] === File Generation Failure Summary ===")
            print("[VERBOSE] No valid .pptx file was found in any message.")
            print(
                "[VERBOSE] Total messages checked: %d"
                % (len(response.messages) if response.messages else 0)
            )
            messages_with_pd = sum(
                1
                for m in (response.messages or [])
                if hasattr(m, "provider_data") and m.provider_data
            )
            print("[VERBOSE] Messages with provider_data: %d" % messages_with_pd)
            print("[VERBOSE] Hint: The agent may not have generated a file. Check if:")
            print(
                "[VERBOSE]   1. The agent's response contained tool use / skill output"
            )
            print("[VERBOSE]   2. The ANTHROPIC_API_KEY has access to the pptx skill")
            print("[VERBOSE]   3. The prompt was understood correctly by the agent")
            print("[VERBOSE] === End Failure Summary ===\n")
        return StepOutput(
            content="Error: No presentation file was generated.",
            success=False,
            stop=True,
        )

    # Extract slide content
    print("\nExtracting slide content...")
    generated_prs = Presentation(generated_file)
    slides_data = []

    # Check template layouts for picture placeholders so the image planner
    # knows which slides have dedicated picture slots in the template.
    template_prs = None
    try:
        template_prs = Presentation(template_path)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))

    total_gen_slides = len(list(generated_prs.slides))

    for idx, slide in enumerate(generated_prs.slides):
        content = _extract_slide_content(slide)

        # Check if the template layout for this slide position has image placeholders
        has_template_image_ph = False
        if template_prs is not None:
            try:
                layout = _find_best_layout(template_prs, idx, total_gen_slides)
                for ph in layout.placeholders:
                    ph_fmt = ph.placeholder_format
                    if ph_fmt is not None and ph_fmt.type == PP_PLACEHOLDER.PICTURE:
                        has_template_image_ph = True
                        break
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))

        slide_info = {
            "index": idx,
            "title": content.title,
            "subtitle": content.subtitle,
            "bullet_count": len(content.body_paragraphs),
            "has_table": len(content.tables) > 0,
            "has_chart": len(content.charts) > 0,
            "has_image": len(content.images) > 0,
            "has_shapes": len(content.shapes_xml) > 0,
            "has_image_placeholder": has_template_image_ph,
        }
        slides_data.append(slide_info)
        print(
            "  Slide %d: '%s' | tables:%d charts:%d images:%d img_ph:%s"
            % (
                idx + 1,
                content.title[:40] if content.title else "",
                len(content.tables),
                len(content.charts),
                len(content.images),
                has_template_image_ph,
            )
        )

    # Store in session state
    session_state["generated_file"] = generated_file
    session_state["slides_data"] = slides_data
    session_state["total_slides"] = len(slides_data)

    slides_summary = json.dumps(slides_data, indent=2)
    return StepOutput(
        content=slides_summary,
        success=True,
    )


# ---------------------------------------------------------------------------
# Step 2: Image Planning (Agent with output_schema)
# ---------------------------------------------------------------------------

# This agent decides which slides need AI-generated images
image_planner = Agent(
    name="Image Planner",
    model=Gemini(
        id="models/gemini-2.5-flash-image",
        response_modalities=["Text", "Image"],
    ),
    instructions=[
        "You are an image planning specialist for PowerPoint presentations.",
        "You will receive a JSON description of slides in a presentation.",
        "For each slide, decide whether an AI-generated image would enhance it.",
        "",
        "GUIDELINES for image decisions:",
        "- Title slides: Usually YES - a hero/background image enhances first impressions.",
        "- Data slides (with tables/charts): Usually NO - data visuals are sufficient.",
        "- Content slides with bullet points: Consider YES if the topic is visual.",
        "- Closing slides: Usually NO - keep the focus on next steps/action items.",
        "- Slides that already have images: ALWAYS NO - they already have visuals.",
        "- Slides with image placeholders (has_image_placeholder=true): STRONGLY YES -",
        "  the template layout has a dedicated picture placeholder that should be filled.",
        "  Generating an image for these slides ensures the placeholder is used properly",
        "  rather than appearing as an empty box.",
        "",
        "When providing image prompts:",
        "- Describe professional, clean, modern illustrations.",
        "- Use abstract or metaphorical imagery, not literal depictions of text.",
        "- Specify the style: 'minimalist', 'corporate', 'flat design', etc.",
        "- Keep prompts under 100 words.",
        "- Make images suitable for a professional business presentation.",
    ],
    output_schema=ImagePlan,
    markdown=False,
)


# ---------------------------------------------------------------------------
# Step 3: Image Generation
# ---------------------------------------------------------------------------


def step_generate_images(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Generate images for slides that need them using NanoBanana (powered by Gemini).

    This step:
    1. Reads the image plan from step 2 (produced by the Gemini image planner)
    2. For each slide marked as needing an image, generates one with NanoBanana
    3. Stores the generated image bytes in session_state keyed by slide index
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)

    print("\n" + "=" * 60)
    print("Step 3: Generating images with NanoBanana...")
    print("=" * 60)

    # Parse the image plan from step 2
    image_plan_content = step_input.previous_step_content
    generated_images: Dict[int, bytes] = {}

    if not image_plan_content:
        print("No image plan received. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(content="No images to generate.", success=True)

    # Parse the plan
    try:
        if isinstance(image_plan_content, str):
            plan_data = json.loads(image_plan_content)
        elif isinstance(image_plan_content, dict):
            plan_data = image_plan_content
        elif isinstance(image_plan_content, ImagePlan):
            plan_data = image_plan_content.model_dump()
        elif isinstance(image_plan_content, BaseModel):
            plan_data = image_plan_content.model_dump()
        else:
            plan_data = json.loads(str(image_plan_content))
    except (json.JSONDecodeError, TypeError):
        print("Could not parse image plan. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(content="Could not parse image plan.", success=True)

    decisions = plan_data.get("decisions", [])
    slides_needing_images = [d for d in decisions if d.get("needs_image", False)]

    if not slides_needing_images:
        print("Image planner decided no slides need images.")
        session_state["generated_images"] = {}
        return StepOutput(content="No slides need images.", success=True)

    # Check which slides already have images from Claude
    slides_data = session_state.get("slides_data", [])
    slides_with_existing_images = {
        s["index"] for s in slides_data if s.get("has_image", False)
    }

    # Initialize NanoBanana
    google_api_key = os.getenv("GOOGLE_API_KEY")
    if not google_api_key:
        print("GOOGLE_API_KEY not set. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(
            content="GOOGLE_API_KEY not set. Skipped image generation.",
            success=True,
        )

    nano_banana = NanoBananaTools(
        api_key=google_api_key,
        aspect_ratio="16:9",  # Widescreen for presentations
    )

    for decision in slides_needing_images:
        slide_idx = decision["slide_index"]
        prompt = decision.get("image_prompt", "")

        if slide_idx in slides_with_existing_images:
            print(
                "  Slide %d: Already has image from Claude, skipping." % (slide_idx + 1)
            )
            continue

        if not prompt:
            print("  Slide %d: No prompt provided, skipping." % (slide_idx + 1))
            continue

        print("  Slide %d: Generating image..." % (slide_idx + 1))
        print("    Prompt: %s" % prompt[:80])

        try:
            result = nano_banana.create_image(prompt)
            if hasattr(result, "images") and result.images:
                for img in result.images:
                    if hasattr(img, "content") and img.content:
                        generated_images[slide_idx] = img.content
                        print(
                            "    Generated successfully (%d bytes)" % len(img.content)
                        )
                        break
            elif isinstance(result, str) and "successfully" in result.lower():
                print("    Image generated but no bytes returned.")
        except Exception as e:
            print("    Failed to generate image: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()

    session_state["generated_images"] = generated_images
    return StepOutput(
        content="Generated %d image(s) for presentation." % len(generated_images),
        success=True,
    )


# ---------------------------------------------------------------------------
# Step 4: Template Assembly
# ---------------------------------------------------------------------------


def step_assemble_template(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Apply template styling and assemble the final presentation.

    Handles image placeholder detection on each template layout so that
    generated images are inserted into dedicated picture placeholders when
    available, rather than being added as free-floating pictures.

    This step:
    1. Opens the template and generated presentation
    2. Creates slides from template layouts with extracted content
    3. Inserts AI-generated images into picture placeholders (or as free-floating)
    4. Saves the final output
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)

    print("\n" + "=" * 60)
    print("Step 4: Assembling final presentation with template...")
    print("=" * 60)

    template_path = session_state.get("template_path", "")
    generated_file = session_state.get("generated_file", "")
    output_path = session_state.get("output_path", "presentation_from_template.pptx")
    generated_images = session_state.get("generated_images", {})

    if not generated_file or not os.path.isfile(generated_file):
        return StepOutput(
            content="Error: Generated file not found.",
            success=False,
            stop=True,
        )

    print("Template: %s" % template_path)
    print("Generated: %s" % generated_file)
    print("Output: %s" % output_path)

    # Open the generated presentation
    generated_prs = Presentation(generated_file)
    generated_slides = list(generated_prs.slides)
    total_slides = len(generated_slides)

    if VERBOSE:
        print("[VERBOSE] Generated presentation has %d slides" % total_slides)

    if total_slides == 0:
        shutil.copy2(template_path, output_path)
        return StepOutput(
            content="Warning: No slides found. Copied template as-is.",
            success=True,
        )

    # Create output from template
    shutil.copy2(template_path, output_path)
    output_prs = Presentation(output_path)
    slide_width = output_prs.slide_width
    slide_height = output_prs.slide_height

    if VERBOSE:
        print(
            "[VERBOSE] Template layouts available: %s"
            % [layout.name for layout in output_prs.slide_layouts]
        )

    # Remove existing slides from template
    while len(output_prs.slides._sldIdLst) > 0:
        sldId = output_prs.slides._sldIdLst[0]
        rId = sldId.get(
            etree.QName(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                "id",
            )
        )
        if rId is not None:
            output_prs.part.drop_rel(rId)
        output_prs.slides._sldIdLst.remove(sldId)

    print("Cleared template slides. Building final presentation...")

    # For each generated slide, create a template-styled slide
    for idx, gen_slide in enumerate(generated_slides):
        content = _extract_slide_content(gen_slide)

        # Look up generated image for this slide.
        # The keys may be int or str depending on how they were stored.
        gen_img = generated_images.get(idx) or generated_images.get(str(idx))

        # Find layout FIRST so we can detect picture placeholders before
        # deciding whether to add a free-floating image.
        layout = _find_best_layout(output_prs, idx, total_slides)

        # Detect picture placeholders on the chosen template layout.
        # Uses both PP_PLACEHOLDER enum and XML-level fallback detection.
        for ph in layout.placeholders:
            ph_fmt = ph.placeholder_format
            if ph_fmt is not None:
                is_pic_ph = False
                # Strategy 1: int comparison with raw OOXML value 18
                try:
                    if ph_fmt.type is not None and (
                        int(ph_fmt.type) == 18 or str(ph_fmt.type) == "PICTURE (18)"
                    ):
                        is_pic_ph = True
                except (ValueError, TypeError):
                    pass
                # Strategy 2: PP_PLACEHOLDER enum
                if not is_pic_ph:
                    try:
                        if ph_fmt.type == PP_PLACEHOLDER.PICTURE:
                            is_pic_ph = True
                    except Exception as e:
                        if VERBOSE:
                            print("[VERBOSE] Exception suppressed: %s" % str(e))
                # Strategy 3: XML-level fallback
                if not is_pic_ph:
                    try:
                        nsmap = {
                            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                        }
                        ph_elem = ph._element.find(".//p:ph", nsmap)
                        if ph_elem is not None and ph_elem.get("type") in (
                            "pic",
                            "clipArt",
                        ):
                            is_pic_ph = True
                    except Exception as e:
                        if VERBOSE:
                            print("[VERBOSE] Exception suppressed: %s" % str(e))
                if is_pic_ph:
                    content.has_image_placeholder = True
                    if ph_fmt.idx not in content.image_placeholder_indices:
                        content.image_placeholder_indices.append(ph_fmt.idx)

        # Only add as free-floating picture if the layout has NO picture
        # placeholders. When picture placeholders exist, the image will be
        # inserted into the placeholder by _populate_slide instead.
        if gen_img is not None and not content.has_image_placeholder:
            content.images.append(
                ImageData(
                    blob=gen_img,
                    left=0,
                    top=0,
                    width=int(Inches(8.0)),
                    height=int(Inches(4.5)),
                    content_type="image/png",
                )
            )

        visual_info = []
        if content.tables:
            visual_info.append("%d table(s)" % len(content.tables))
        if content.images:
            visual_info.append("%d image(s)" % len(content.images))
        if content.charts:
            visual_info.append("%d chart(s)" % len(content.charts))
        if content.has_image_placeholder:
            visual_info.append("img placeholder(s)")
        visual_str = ", ".join(visual_info) if visual_info else "text only"

        print(
            "  Slide %d: layout '%s' | title: '%s' | %s"
            % (
                idx + 1,
                layout.name,
                content.title[:40] if content.title else "",
                visual_str,
            )
        )

        new_slide = output_prs.slides.add_slide(layout)
        # Pass generated image bytes to _populate_slide for insertion
        # into picture placeholders
        _populate_slide(
            new_slide,
            content,
            slide_width,
            slide_height,
            generated_image_bytes=gen_img,
        )

    output_prs.save(output_path)
    print("\nSaved final presentation: %s" % output_path)

    return StepOutput(
        content="Presentation saved to %s (%d slides)" % (output_path, total_slides),
        success=True,
    )


# ---------------------------------------------------------------------------
# CLI and Main
# ---------------------------------------------------------------------------


def parse_args():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description="Workflow: Generate PowerPoint with Claude, add AI images, apply template."
    )
    parser.add_argument(
        "--template",
        "-t",
        required=True,
        help="Path to the .pptx template file.",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="presentation_from_template.pptx",
        help="Output filename (default: presentation_from_template.pptx).",
    )
    parser.add_argument(
        "--prompt",
        "-p",
        default=None,
        help="Custom prompt for the presentation content.",
    )
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Skip AI image generation (Steps 2 and 3).",
    )
    parser.add_argument(
        "--no-stream",
        action="store_true",
        help="Disable streaming mode for Claude agent (more reliable for shorter prompts, "
        "but may timeout on complex presentations).",
    )
    parser.add_argument(
        "--verbose",
        "-v",
        action="store_true",
        help="Enable verbose/debug logging for troubleshooting.",
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()

    VERBOSE = args.verbose

    if not os.path.isfile(args.template):
        print("Error: Template file not found: %s" % args.template)
        sys.exit(1)

    if not args.template.endswith(".pptx"):
        print("Error: Template file must be a .pptx file")
        sys.exit(1)

    if not os.getenv("ANTHROPIC_API_KEY"):
        raise ValueError("ANTHROPIC_API_KEY environment variable not set")

    # Setup output directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(script_dir, "output")
    os.makedirs(output_dir, exist_ok=True)

    output_path = args.output
    if not os.path.isabs(output_path):
        output_path = os.path.join(output_dir, output_path)

    # Default prompt
    prompt = args.prompt
    if not prompt:
        prompt = (
            "Create a 6-slide business presentation with the following structure.\n"
            "Rules: Use one clear title per slide, 4-6 concise bullet points for text slides,\n"
            "and include tables or charts where specified. Do not apply custom styling.\n"
            "\n"
            "Slide 1 - Title Slide:\n"
            "  Title: 'Strategic Overview 2026'\n"
            "  Subtitle: 'Annual Business Review and Forward Plan'\n"
            "\n"
            "Slide 2 - Market Analysis (text with bullets):\n"
            "  Title: 'Market Analysis'\n"
            "  Bullets: Key market trends, growth opportunities, competitive landscape,\n"
            "  emerging technologies, and regulatory changes.\n"
            "\n"
            "Slide 3 - Financial Performance (include a TABLE):\n"
            "  Title: 'Financial Performance'\n"
            "  Table with columns: Metric, Q1, Q2, Q3, Q4\n"
            "  Rows: Revenue ($M), Costs ($M), Profit ($M), Growth (%)\n"
            "\n"
            "Slide 4 - Revenue Trends (include a BAR CHART):\n"
            "  Title: 'Quarterly Revenue Trends'\n"
            "  Bar chart showing quarterly revenue for 2024 vs 2025.\n"
            "\n"
            "Slide 5 - Our Strategy (text with bullets):\n"
            "  Title: 'Strategic Priorities'\n"
            "  Bullets: Three-pillar growth approach, market expansion,\n"
            "  product innovation, operational excellence.\n"
            "\n"
            "Slide 6 - Closing Slide:\n"
            "  Title: 'Next Steps'\n"
            "  Bullets: Implementation timeline, key milestones, success metrics.\n"
            "\n"
            "Save as 'generated_content.pptx'"
        )

    print("=" * 60)
    print("PowerPoint Template Workflow")
    print("=" * 60)
    print("Template: %s" % args.template)
    print("Output:   %s" % output_path)
    print("Images:   %s" % ("disabled" if args.no_images else "enabled"))
    print("Stream:   %s" % ("disabled" if args.no_stream else "enabled"))
    if VERBOSE:
        print("Verbose:  enabled")

    # Build workflow steps
    steps: List[Step] = [
        Step(name="Content Generation", executor=step_generate_content),
    ]

    if not args.no_images:
        steps.append(Step(name="Image Planning", agent=image_planner))
        steps.append(Step(name="Image Generation", executor=step_generate_images))

    steps.append(Step(name="Template Assembly", executor=step_assemble_template))

    # Create and run the workflow
    #
    # session_state schema:
    #   template_path      (str)  - Path to the .pptx template file
    #   output_dir         (str)  - Directory for intermediate output files
    #   output_path        (str)  - Path for the final assembled presentation
    #   generated_file     (str)  - Path to the Claude-generated .pptx file (set by step 1)
    #   slides_data        (list) - Per-slide metadata dicts (set by step 1)
    #   total_slides       (int)  - Number of slides in the generated presentation
    #   generated_images   (dict) - Mapping of slide index -> image bytes (set by step 3)
    #   verbose            (bool) - Whether verbose logging is enabled
    #   stream             (bool) - Whether to use streaming mode for Claude agent
    workflow = Workflow(
        name="PowerPoint Template Workflow",
        steps=steps,
        session_state={
            "template_path": args.template,
            "output_dir": output_dir,
            "output_path": output_path,
            "generated_file": "",
            "slides_data": [],
            "total_slides": 0,
            "generated_images": {},
            "verbose": args.verbose,
            "stream": not args.no_stream,
        },
    )

    workflow.print_response(input=prompt, markdown=True)

    print("\n" + "=" * 60)
    print("Workflow complete!")
    print("Output: %s" % output_path)
    print("=" * 60)
