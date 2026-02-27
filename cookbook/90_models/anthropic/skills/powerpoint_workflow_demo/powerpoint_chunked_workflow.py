"""
Agno Workflow: Chunked PowerPoint Generation Pipeline.

A chunked workflow that overcomes Claude API limitations for large presentations
by splitting generation into manageable chunks, then merging the results.

Problem solved: Single Claude API calls fail for 10+ slide presentations.
Solution: Generate slides in configurable chunks (default: 3 slides per call),
          then merge all chunks into one final presentation.

Architecture:
  This file is a thin orchestration layer built on top of powerpoint_template_workflow.py.
  It imports all helpers, agents, Pydantic models, and step functions from that file via
  a wildcard import, then adds the chunked orchestration logic on top.

  powerpoint_template_workflow.py  — Core pipeline: content gen, images, template assembly,
                                     visual review, all helper functions (~7000 lines)
  powerpoint_chunked_workflow.py   — Chunked orchestration layer (~1500 lines, this file)

Relationship between the two files:
  - powerpoint_template_workflow.py is self-contained and can run standalone (single Claude
    API call, suitable for short presentations of up to ~7 slides).
  - powerpoint_chunked_workflow.py wraps the same template/image/review logic via wildcard
    import so that large presentations (8-15+ slides) are split into chunks and merged.
  - Do NOT modify powerpoint_template_workflow.py to add chunking logic; keep them separate.

Workflow steps:
  Step 1  Optimize & Plan    - LLM analyzes prompt, decides slide count, creates storyboard
  Step 2  Generate Chunks    - Call Claude pptx skill for each chunk of slides
  Step 3  Process Chunks     - Apply template + image pipeline per chunk (if template provided)
  Step 4  Visual Review      - Optional per-chunk visual QA (if --visual-review + template)
  Step 5  Merge Chunks       - Merge all processed chunks into the final PPTX

Prerequisites:
- uv pip install agno anthropic python-pptx google-genai pillow
- export ANTHROPIC_API_KEY="your_api_key_here"
- export GOOGLE_API_KEY="your_google_api_key_here" (for image generation)
- A .pptx template file (optional)

Usage:
    # Basic usage (auto-decide slide count, 3 slides per chunk):
    .venvs/demo/bin/python powerpoint_chunked_workflow.py \\
        -p "Create a presentation about AI in healthcare"

    # With template, 4 slides per chunk:
    .venvs/demo/bin/python powerpoint_chunked_workflow.py \\
        -t my_template.pptx --chunk-size 4

    # Large presentation with visual review (5 passes max):
    .venvs/demo/bin/python powerpoint_chunked_workflow.py \\
        -t my_template.pptx -p "12-slide enterprise AI strategy deck" \\
        --chunk-size 3 --visual-review --visual-passes 5

    # Quick generation without images or template:
    .venvs/demo/bin/python powerpoint_chunked_workflow.py \\
        -p "Startup pitch deck for SaaS product" --no-images

CLI Flags:
    --template, -t       Path to .pptx template (optional). Without it, skips
                         template assembly and visual review; just merges raw chunks.
    --output, -o         Output filename (default: presentation_chunked.pptx).
    --prompt, -p         User prompt describing the presentation.
    --no-images          Skip AI image generation.
    --no-stream          Disable streaming mode for Claude agent.
    --min-images         Minimum slides that must have images (default: 1).
    --visual-review      Enable visual QA with Gemini vision per chunk.
    --footer-text        Footer text for all slides.
    --date-text          Date text for footer date placeholder.
    --show-slide-numbers Preserve slide number placeholder on all slides.
    --verbose, -v        Enable verbose/debug logging.
    --chunk-size         Number of slides per Claude API chunk call (default: 3).
    --max-retries        Max retries per chunk on failure (default: 2).
    --visual-passes      Maximum visual inspection passes per chunk (default: 3).

Logging conventions:
    Always printed:
        [STEP_NAME] Starting / result messages
        [TIMING] step_XXX completed in X.Xs
        [ERROR] ...
        [WARNING] ...
        [VISUAL REVIEW MISSING FIX] ...  (always, per spec)
    Verbose-only (requires --verbose / -v):
        [VERBOSE] detailed debug information
"""

import argparse
import copy
import json
import os
import shutil
import sys
import time
import traceback
import zipfile
from typing import Dict, List, Optional

from agno.run.agent import RunOutput

# === WILDCARD IMPORT: Reuse all helpers, agents, models, and step functions ===
# This gives us access to all ~6500 lines of helper logic without duplication.
# Specifically imports: SlideImageDecision, ImagePlan, ShapeIssue, SlideQualityReport,
# PresentationQualityReport, all dataclasses, image_planner, slide_quality_reviewer,
# step_plan_images, step_generate_images, step_assemble_template, step_visual_quality_review,
# and all _helper functions.
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from powerpoint_template_workflow import *  # noqa: F401, F403, E402

from agno.agent import Agent
from agno.models.anthropic import Claude
from agno.workflow.step import Step
from agno.workflow.types import StepInput, StepOutput
from agno.workflow.workflow import Workflow
from anthropic import Anthropic
from file_download_helper import download_skill_files
from pptx import Presentation
from pydantic import BaseModel, Field


# === NEW PYDANTIC MODELS FOR CHUNKED WORKFLOW ===


class SlideStoryboard(BaseModel):
    """Storyboard entry for a single slide in the presentation."""

    slide_number: int = Field(..., description="1-based slide number")
    slide_title: str = Field(..., description="Title for this slide")
    slide_type: str = Field(
        ..., description="Type: title/agenda/content/data/closing"
    )
    key_points: List[str] = Field(..., description="3-5 bullet points for this slide")
    visual_suggestion: str = Field(
        ..., description="Visual element recommendation"
    )
    transition_note: str = Field(
        ..., description="How this slide connects to next"
    )


class StoryboardPlan(BaseModel):
    """Complete storyboard plan for the presentation, produced by the query optimizer."""

    total_slides: int = Field(..., description="Total number of slides")
    presentation_title: str = Field(..., description="Main presentation title")
    target_audience: str = Field(..., description="Target audience")
    tone: str = Field(..., description="e.g. professional, inspiring, technical")
    brand_voice: str = Field(
        ..., description="e.g. authoritative, conversational, data-driven"
    )
    global_context: str = Field(
        ...,
        description="Shared context applicable to all slides - company, product, theme",
    )
    slides: List[SlideStoryboard] = Field(
        ..., description="Per-slide storyboard entries"
    )


# === MODULE-LEVEL AGENTS ===
# Do NOT create agents in loops — define them here at module level.

query_optimizer = Agent(
    name="Presentation Strategist",
    model=Claude(id="claude-sonnet-4-6"),
    betas=["context-1m-2025-08-07"],
    description=(
        "You are a presentation strategist who first searches the web for current, "
        "relevant facts and data about the topic, then creates an optimized presentation "
        "plan with a per-slide storyboard grounded in that research."
    ),
    tools=[
        {
            "type": "web_search_20250305",
            "name": "web_search",
            "max_uses": 5,
        }
    ],
    thinking={"type": "disabled"},
    output_schema=StoryboardPlan,
)


# === HELPER: STORYBOARD MARKDOWN FORMATTING ===


def _format_slide_markdown(slide: SlideStoryboard) -> str:
    """Format a SlideStoryboard as a markdown file for the pptx agent.

    Excludes transition_note (planning meta-info, not useful for content generation)
    to keep context size lean. Includes type, key points, and visual suggestion only.
    """
    points = "\n".join("- %s" % p for p in slide.key_points)
    return (
        "# Slide %d: %s\n\n"
        "**Type:** %s\n\n"
        "## Content\n%s\n\n"
        "**Visual:** %s\n"
    ) % (
        slide.slide_number,
        slide.slide_title,
        slide.slide_type,
        points,
        slide.visual_suggestion,
    )


def _format_global_context_markdown(plan: StoryboardPlan) -> str:
    """Format the global context as a markdown file for the pptx agent.

    Kept concise: title, audience, tone, brand voice, and the 2-3 sentence global context.
    This file is included in every chunk prompt so brevity matters for context size.
    """
    return (
        "# Presentation: %s\n\n"
        "Audience: %s | Tone: %s | Brand Voice: %s\n\n"
        "## Context\n%s\n"
    ) % (
        plan.presentation_title,
        plan.target_audience,
        plan.tone,
        plan.brand_voice,
        plan.global_context,
    )


# === HELPER: EXTRACT SLIDES DATA FROM A CHUNK PPTX ===


def _extract_chunk_slides_data(chunk_file: str) -> List[dict]:
    """Extract basic slide metadata from a PPTX chunk file.

    Returns a list of dicts compatible with the session_state['slides_data'] format
    used by step_plan_images and step_generate_images.
    """
    slides_data = []
    try:
        prs = Presentation(chunk_file)
        for idx, slide in enumerate(prs.slides):
            slide_info: dict = {
                "index": idx,
                "title": "",
                "body": "",
                "has_table": False,
                "has_chart": False,
                "has_image": False,
            }
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for ph_attr in ["placeholder_format"]:
                        ph = getattr(shape, ph_attr, None)
                        if ph and hasattr(ph, "idx") and ph.idx == 0:
                            slide_info["title"] = shape.text_frame.text.strip()
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    slide_info["has_table"] = True
                if shape.shape_type == 3:  # CHART
                    slide_info["has_chart"] = True
            slides_data.append(slide_info)
    except Exception as e:
        print("[WARNING] Could not extract slides data from %s: %s" % (chunk_file, e))
    return slides_data


# === WORKFLOW STEP 1: OPTIMIZE AND PLAN ===


def step_optimize_and_plan(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 1: Enhance the user prompt, decide slide count, and generate a per-slide storyboard.

    Uses the query_optimizer agent (Claude Sonnet) to produce a StoryboardPlan with:
    - Optimal slide count (respects user-specified count, otherwise picks 8-15)
    - Global context applicable to all slides
    - Per-slide storyboard with title, type, key points, visual suggestions
    - Presentation tone and brand voice

    Saves storyboard to individual markdown files in {output_dir}/storyboard/.
    """
    step_start = time.time()

    user_prompt = session_state.get("user_prompt", "")
    output_dir = session_state.get("output_dir", ".")
    chunk_size = session_state.get("chunk_size", 3)
    max_retries = session_state.get("max_retries", 2)

    print("=" * 60)
    print("Step 1: Optimizing query and generating storyboard...")
    print("=" * 60)
    print("User prompt: %s" % user_prompt[:200])

    storyboard_dir = os.path.join(output_dir, "storyboard")
    os.makedirs(storyboard_dir, exist_ok=True)

    optimizer_prompt = (
        "Analyze the following user request for a PowerPoint presentation and create an optimized storyboard.\n\n"
        "USER REQUEST:\n%s\n\n"
        "STEP 1 — RESEARCH FIRST:\n"
        "Before planning slides, use web_search to find 2-4 relevant facts, statistics, or examples "
        "for the topic. Search for current data (e.g. market size, adoption rates, key trends). "
        "Use these findings to ground the storyboard in real, specific information.\n\n"
        "STEP 2 — BUILD THE STORYBOARD:\n"
        "1. If the user specifies a slide count (e.g. '12 slides', '10-slide deck'), honor it exactly.\n"
        "2. If not specified, decide the optimal count: typically 8-15 slides for professional decks.\n"
        "   - Simple topics: 8-10 slides. Complex/technical topics: 12-15 slides.\n"
        "   - Do NOT add unnecessary slides; quality over quantity.\n"
        "3. Define a clear tone and brand voice appropriate to the topic.\n"
        "4. Write global_context as 2-3 focused sentences covering: the core topic, target audience, "
        "and central theme or key message. Include a specific fact or statistic from your research.\n"
        "5. For each slide, provide:\n"
        "   - A concise, descriptive title (5-8 words max)\n"
        "   - slide_type: one of title, agenda, content, data, closing\n"
        "   - 3-4 key_points: each a single sentence (10-20 words), specific and actionable. "
        "     Include real data or examples where relevant. Avoid vague filler bullets.\n"
        "   - visual_suggestion: one concrete line — specify chart type + data, image concept, or 'none'. "
        "     Example: 'bar chart: AI adoption rate by industry 2023' not just 'chart'.\n"
        "   - transition_note: one brief sentence connecting this slide to the next.\n"
        "6. Ensure continuity: the storyboard should feel like a coherent narrative arc.\n"
        "7. Use professional language. Do not add emojis or overly casual language.\n"
    ) % user_prompt

    if VERBOSE:  # noqa: F405
        print(
            "[VERBOSE] Full optimizer prompt (%d chars):\n%s"
            % (len(optimizer_prompt), optimizer_prompt)
        )

    try:
        response = query_optimizer.run(optimizer_prompt, stream=False)
    except Exception as e:
        print("[ERROR] Query optimizer failed: %s" % str(e))
        traceback.print_exc()
        return StepOutput(
            content="Query optimization failed: %s" % str(e), success=False
        )

    # Parse the StoryboardPlan from response
    plan: Optional[StoryboardPlan] = None
    if response and response.content:
        content = response.content
        if isinstance(content, StoryboardPlan):
            plan = content
        elif isinstance(content, dict):
            try:
                plan = StoryboardPlan(**content)
            except Exception as e:
                print("[ERROR] Failed to parse StoryboardPlan from dict: %s" % e)
        elif isinstance(content, str):
            try:
                plan = StoryboardPlan.model_validate_json(content)
            except Exception as e:
                print("[ERROR] Failed to parse StoryboardPlan from JSON string: %s" % e)

    if not plan:
        print("[ERROR] No valid storyboard plan produced.")
        return StepOutput(content="No storyboard plan produced.", success=False)

    print(
        "Storyboard plan: '%s' (%d slides, tone: %s)"
        % (plan.presentation_title, plan.total_slides, plan.tone)
    )

    if VERBOSE:  # noqa: F405
        print("[VERBOSE] Full storyboard JSON:\n%s" % plan.model_dump_json(indent=2))

    # Save global context markdown
    global_context_path = os.path.join(storyboard_dir, "global_context.md")
    with open(global_context_path, "w", encoding="utf-8") as f:
        f.write(_format_global_context_markdown(plan))
    print("Saved global context: %s" % global_context_path)

    # Save per-slide storyboard markdown files
    for slide in plan.slides:
        slide_path = os.path.join(storyboard_dir, "slide_%03d.md" % slide.slide_number)
        slide_md = _format_slide_markdown(slide)
        if VERBOSE:  # noqa: F405
            print(
                "[VERBOSE] Slide %d storyboard:\n%s"
                % (slide.slide_number, slide_md)
            )
        with open(slide_path, "w", encoding="utf-8") as f:
            f.write(slide_md)
    print("Saved %d slide storyboard files to: %s" % (len(plan.slides), storyboard_dir))

    # Store in session_state
    session_state["storyboard"] = plan
    session_state["total_slides"] = plan.total_slides
    session_state["storyboard_dir"] = storyboard_dir
    session_state["chunk_size"] = chunk_size
    session_state["max_retries"] = max_retries

    step_elapsed = time.time() - step_start
    print("[TIMING] step_optimize_and_plan completed in %.1fs" % step_elapsed)

    summary = (
        "Storyboard created: '%s' | %d slides | tone: %s | brand voice: %s | chunk size: %d | Duration: %.1fs"
    ) % (plan.presentation_title, plan.total_slides, plan.tone, plan.brand_voice, chunk_size, step_elapsed)
    return StepOutput(content=summary, success=True)


# === HELPER: GENERATE A SINGLE CHUNK VIA CLAUDE PPTX SKILL ===


def generate_chunk_pptx(
    chunk_slides: List[SlideStoryboard],
    session_state: Dict,
    chunk_idx: int,
) -> Optional[str]:
    """Call the Claude pptx skill for a chunk of slides with retry logic.

    Creates a fresh agent per call (not reused across chunks) and applies
    exponential backoff on retries.

    Args:
        chunk_slides: List of SlideStoryboard objects for this chunk.
        session_state: Shared workflow session state.
        chunk_idx: 0-based chunk index (used for file naming and logging).

    Returns:
        Path to the generated chunk PPTX file, or None if all attempts failed.
    """
    storyboard: StoryboardPlan = session_state["storyboard"]
    storyboard_dir = session_state["storyboard_dir"]
    output_dir = session_state["output_dir"]
    max_retries = session_state.get("max_retries", 2)

    # Load global context
    global_context_path = os.path.join(storyboard_dir, "global_context.md")
    global_context = ""
    if os.path.exists(global_context_path):
        with open(global_context_path, encoding="utf-8") as f:
            global_context = f.read()

    # Load per-slide markdown for this chunk
    slide_details = []
    for s in chunk_slides:
        md_path = os.path.join(storyboard_dir, "slide_%03d.md" % s.slide_number)
        if os.path.exists(md_path):
            with open(md_path, encoding="utf-8") as f:
                slide_details.append(f.read())
        else:
            # Fallback: format inline
            slide_details.append(_format_slide_markdown(s))

    first_slide = chunk_slides[0].slide_number
    last_slide = chunk_slides[-1].slide_number

    chunk_prompt = (
        "## Global Presentation Context\n"
        "%s\n\n"
        "## Task: Generate slides %d through %d of %d\n\n"
        "You are generating a CHUNK of a larger presentation. "
        "This chunk contains %d slides.\n"
        "Maintain the presentation's tone (%s) and brand voice (%s).\n"
        "These are slides %d-%d of the full %d-slide deck titled \"%s\".\n\n"
        "## Per-Slide Content for This Chunk:\n\n"
        "%s\n\n"
        "Please generate EXACTLY %d slides for this chunk with the content described above.\n"
        "Do not add extra slides. Do not include slide numbers outside the range %d-%d.\n"
        "Use clean formatting without custom fonts or colors. "
        "Include tables or charts only where explicitly suggested.\n"
        "Save the output as 'chunk_%03d.pptx'."
    ) % (
        global_context,
        first_slide,
        last_slide,
        storyboard.total_slides,
        len(chunk_slides),
        storyboard.tone,
        storyboard.brand_voice,
        first_slide,
        last_slide,
        storyboard.total_slides,
        storyboard.presentation_title,
        "\n\n---\n\n".join(slide_details),
        len(chunk_slides),
        first_slide,
        last_slide,
        chunk_idx,
    )

    # if VERBOSE:  # noqa: F405
    #     print(
    #         "[VERBOSE] Chunk %d prompt (%d chars):\n%s"
    #         % (chunk_idx, len(chunk_prompt), chunk_prompt[:1000])
    #     )

    chunk_output_path = os.path.join(output_dir, "chunk_%03d.pptx" % chunk_idx)

    # Create a fresh agent per chunk call — do NOT reuse across calls
    chunk_agent = Agent(
        name="Chunk Generator %d" % chunk_idx,
        model=Claude(
            id="claude-sonnet-4-6",
            betas=["context-1m-2025-08-07"],
            skills=[{"type": "anthropic", "skill_id": "pptx", "version": "latest"}],
        ),
        instructions=[
            "You are a structured content generator for PowerPoint presentations.",
            "Generate EXACTLY the number of slides specified in the task.",
            "Use one clear title per slide with concise bullet points.",
            "Do NOT apply custom fonts, colors, or theme styling.",
            "Do NOT add animations, transitions, or speaker notes.",
            "Keep tables to max 6 rows x 5 columns.",
        ],
        markdown=True,
    )

    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

    for attempt in range(max_retries + 1):
        attempt_start = time.time()

        if attempt > 0:
            delay_ms = int(1000 * (2 ** (attempt - 1)))  # exponential: 1000ms, 2000ms
            print(
                "[CHUNK %d] Retry %d/%d after %dms delay..."
                % (chunk_idx, attempt, max_retries, delay_ms)
            )
            time.sleep(delay_ms / 1000.0)

        print(
            "[CHUNK %d] API call attempt %d/%d (slides %d-%d)..."
            % (chunk_idx, attempt + 1, max_retries + 1, first_slide, last_slide)
        )

        try:
            response = None
            event_count = 0
            for event in chunk_agent.run(chunk_prompt, stream=True, yield_run_output=True):
                event_count += 1
                if isinstance(event, RunOutput):
                    response = event

            if response is None:
                print("[CHUNK %d] No RunOutput received after %d events." % (chunk_idx, event_count))
                attempt_elapsed = time.time() - attempt_start
                print(
                    "[TIMING] Chunk %d attempt %d/%d: %.1fs (no output)"
                    % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
                )
                continue

            if VERBOSE:  # noqa: F405
                msg_count = len(response.messages) if response.messages else 0
                print(
                    "[VERBOSE] Chunk %d attempt %d: received %d events, %d messages"
                    % (chunk_idx, attempt + 1, event_count, msg_count)
                )
                if response.messages:
                    for m_idx, msg in enumerate(response.messages):
                        print(
                            "[VERBOSE] Chunk %d message %d: type=%s role=%s has_provider_data=%s"
                            % (
                                chunk_idx,
                                m_idx,
                                type(msg).__name__,
                                getattr(msg, "role", "N/A"),
                                bool(getattr(msg, "provider_data", None)),
                            )
                        )

        except Exception as e:
            print("[CHUNK %d] Attempt %d/%d failed with error: %s" % (chunk_idx, attempt + 1, max_retries + 1, e))
            attempt_elapsed = time.time() - attempt_start
            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (error)"
                % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
            )
            if attempt == max_retries:
                print(
                    "[CHUNK %d] All %d attempts failed. Skipping chunk."
                    % (chunk_idx, max_retries + 1)
                )
            continue

        # Try to download the generated file from message provider_data
        generated_file = None

        if response.messages:
            for msg in response.messages:
                if hasattr(msg, "provider_data") and msg.provider_data:
                    if VERBOSE:  # noqa: F405
                        print(
                            "[VERBOSE] Chunk %d: attempting file download from message provider_data..."
                            % chunk_idx
                        )
                    try:
                        files = download_skill_files(
                            msg.provider_data, client, output_dir=output_dir
                        )
                    except Exception as e:
                        print(
                            "[CHUNK %d] download_skill_files (message) failed: %s"
                            % (chunk_idx, e)
                        )
                        files = []

                    if VERBOSE:  # noqa: F405
                        print(
                            "[VERBOSE] Chunk %d: download returned files: %s"
                            % (chunk_idx, files)
                        )

                    if files:
                        for f in files:
                            if not f.endswith(".pptx"):
                                continue
                            try:
                                Presentation(f)
                                generated_file = f
                                break
                            except Exception:
                                continue
                    if generated_file:
                        break

        # Fallback: try response.model_provider_data
        if not generated_file and hasattr(response, "model_provider_data") and response.model_provider_data:
            if VERBOSE:  # noqa: F405
                print(
                    "[VERBOSE] Chunk %d: trying fallback model_provider_data download..."
                    % chunk_idx
                )
            try:
                files = download_skill_files(
                    response.model_provider_data, client, output_dir=output_dir
                )
                if VERBOSE:  # noqa: F405
                    print(
                        "[VERBOSE] Chunk %d: fallback download returned files: %s"
                        % (chunk_idx, files)
                    )
                for f in files:
                    if not f.endswith(".pptx"):
                        continue
                    try:
                        Presentation(f)
                        generated_file = f
                        break
                    except Exception:
                        continue
            except Exception as e:
                print(
                    "[CHUNK %d] download_skill_files (fallback) failed: %s" % (chunk_idx, e)
                )

        attempt_elapsed = time.time() - attempt_start

        if generated_file and os.path.exists(generated_file):
            # Normalize to standard chunk name
            if generated_file != chunk_output_path:
                shutil.copy2(generated_file, chunk_output_path)
                generated_file = chunk_output_path
            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (success)"
                % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
            )
            print("[CHUNK %d] Successfully generated: %s" % (chunk_idx, generated_file))
            return generated_file
        else:
            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (no file returned)"
                % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
            )
            print("[CHUNK %d] Attempt %d/%d produced no file." % (chunk_idx, attempt + 1, max_retries + 1))

    print(
        "[CHUNK %d] All %d attempts failed. Skipping chunk."
        % (chunk_idx, max_retries + 1)
    )
    return None


# === WORKFLOW STEP 2: GENERATE CHUNKS ===


def step_generate_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 2: Orchestrate chunked PPTX generation across all slide groups.

    Splits the full storyboard into chunks of {chunk_size} slides, calls the
    Claude pptx skill for each chunk with retry logic, and stores results.

    A 1-second inter-chunk delay is applied between calls to avoid rate limits.
    """
    step_start = time.time()

    storyboard: Optional[StoryboardPlan] = session_state.get("storyboard")
    if not storyboard:
        print("[ERROR] No storyboard found in session_state.")
        return StepOutput(content="No storyboard found.", success=False)

    chunk_size = session_state.get("chunk_size", 3)
    slides = storyboard.slides

    # Build chunk list
    chunks = [slides[i : i + chunk_size] for i in range(0, len(slides), chunk_size)]

    print("\n" + "=" * 60)
    print("Step 2: Generating presentation chunks...")
    print("=" * 60)
    print(
        "Total slides: %d | Chunk size: %d | Number of chunks: %d"
        % (len(slides), chunk_size, len(chunks))
    )

    if VERBOSE:  # noqa: F405
        for ci, chunk in enumerate(chunks):
            slide_nums = [s.slide_number for s in chunk]
            print("[VERBOSE] Chunk %d: slides %s" % (ci, slide_nums))

    chunk_files: List[Optional[str]] = []
    successful = 0
    total_chunks = len(chunks)

    for chunk_idx, chunk_slides in enumerate(chunks):
        chunk_start = time.time()
        print(
            "[GENERATE] Chunk %d/%d: slides %d-%d"
            % (
                chunk_idx + 1,
                total_chunks,
                chunk_slides[0].slide_number,
                chunk_slides[-1].slide_number,
            )
        )

        chunk_file = generate_chunk_pptx(chunk_slides, session_state, chunk_idx)

        chunk_elapsed = time.time() - chunk_start
        if chunk_file:
            print(
                "[TIMING] Chunk %d/%d done in %.1fs -> %s"
                % (chunk_idx + 1, total_chunks, chunk_elapsed, chunk_file)
            )
        else:
            print(
                "[TIMING] Chunk %d/%d FAILED after %.1fs (skipping)"
                % (chunk_idx + 1, total_chunks, chunk_elapsed)
            )
            print(
                "[WARNING] Chunk %d failed (slides %d-%d). Continuing..."
                % (
                    chunk_idx,
                    chunk_slides[0].slide_number,
                    chunk_slides[-1].slide_number,
                )
            )

        chunk_files.append(chunk_file)

        if chunk_file:
            successful += 1

        # Inter-chunk delay (not for last chunk) to avoid rate limits
        if chunk_idx < total_chunks - 1:
            inter_delay = 1.0
            print("[GENERATE] Waiting %.1fs before next chunk..." % inter_delay)
            time.sleep(inter_delay)

    session_state["chunk_files"] = chunk_files
    session_state["chunk_slide_groups"] = chunks

    step_elapsed = time.time() - step_start
    failed = total_chunks - successful
    print(
        "\n[TIMING] step_generate_chunks completed in %.1fs (%d chunks: %d succeeded, %d failed)"
        % (step_elapsed, total_chunks, successful, failed)
    )

    summary = "%d of %d chunks generated successfully. Duration: %.1fs" % (
        successful,
        total_chunks,
        step_elapsed,
    )
    return StepOutput(content=summary, success=successful > 0)


# === WORKFLOW STEP 3: PROCESS CHUNKS (TEMPLATE + IMAGES) ===


def step_process_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 3: Apply template assembly and image pipeline to each chunk.

    For each successfully generated chunk, runs:
    1. Image planning (which slides need AI-generated images)
    2. Image generation (NanoBanana)
    3. Template assembly (if --template is provided)

    Each chunk is processed with a temporary session_state copy that adapts
    the existing step functions to work on individual chunk files.
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 3: Processing chunks (images + template assembly)...")
    print("=" * 60)

    chunk_files: List[Optional[str]] = session_state.get("chunk_files", [])
    chunk_slide_groups: List[List[SlideStoryboard]] = session_state.get(
        "chunk_slide_groups", []
    )
    template_path = session_state.get("template_path", "")
    output_dir = session_state.get("output_dir", ".")
    no_images = session_state.get("no_images", False)

    processed_chunks: Dict[int, Optional[str]] = {}
    total_process_chunks = len(chunk_files)

    for chunk_idx, chunk_file in enumerate(chunk_files):
        chunk_proc_start = time.time()

        if chunk_file is None:
            print("[PROCESS] Chunk %d (%d/%d): skipped (no file)." % (chunk_idx, chunk_idx + 1, total_process_chunks))
            processed_chunks[chunk_idx] = None
            continue

        print(
            "\n[PROCESS] Chunk %d (%d/%d): processing %s"
            % (chunk_idx, chunk_idx + 1, total_process_chunks, chunk_file)
        )

        # Determine which slides are in this chunk
        chunk_slides = chunk_slide_groups[chunk_idx] if chunk_idx < len(chunk_slide_groups) else []
        slides_data = _extract_chunk_slides_data(chunk_file)
        total_chunk_slides = len(slides_data)

        assembled_path = os.path.join(
            output_dir, "chunk_%03d_assembled.pptx" % chunk_idx
        )

        # Build a temporary session_state for the existing step functions
        chunk_session = dict(session_state)
        chunk_session["generated_file"] = chunk_file
        chunk_session["total_slides"] = total_chunk_slides
        chunk_session["slides_data"] = slides_data
        chunk_session["output_path"] = assembled_path
        chunk_session["generated_images"] = {}
        # Use chunk-specific output subdirectory to avoid collisions
        chunk_out_subdir = os.path.join(output_dir, "chunk_%03d_work" % chunk_idx)
        os.makedirs(chunk_out_subdir, exist_ok=True)
        chunk_session["output_dir"] = chunk_out_subdir

        if VERBOSE:  # noqa: F405
            print(
                "[VERBOSE] Chunk %d session state keys: %s"
                % (chunk_idx, sorted(chunk_session.keys()))
            )

        # Adjust src_slide dimensions if not set
        if not chunk_session.get("src_slide_width"):
            try:
                prs = Presentation(chunk_file)
                chunk_session["src_slide_width"] = prs.slide_width
                chunk_session["src_slide_height"] = prs.slide_height
            except Exception:
                pass

        current_file = chunk_file

        # --- Image planning ---
        if not no_images:
            print("[PROCESS] Chunk %d: running image planning..." % chunk_idx)
            try:
                # Build slides JSON for image planner
                slides_json = json.dumps(slides_data, indent=2)
                user_prompt = session_state.get("user_prompt", "professional presentation")
                combined_message = (
                    'Presentation topic: "%s"\n\nSlide metadata:\n%s\n\n'
                    "Analyze each slide and decide which ones need AI-generated images.\n"
                    "Consider the presentation topic when writing image prompts."
                ) % (user_prompt, slides_json)

                img_plan_response = image_planner.run(combined_message, stream=False)  # noqa: F405

                if img_plan_response and img_plan_response.content:
                    content = img_plan_response.content

                    if VERBOSE:  # noqa: F405
                        if isinstance(content, BaseModel):
                            print(
                                "[VERBOSE] Chunk %d image plan:\n%s"
                                % (chunk_idx, content.model_dump_json(indent=2))
                            )
                        else:
                            print(
                                "[VERBOSE] Chunk %d image plan content: %s"
                                % (chunk_idx, str(content)[:500])
                            )

                    if isinstance(content, BaseModel):
                        plan_json = content.model_dump_json()
                    elif isinstance(content, dict):
                        plan_json = json.dumps(content)
                    else:
                        plan_json = str(content)

                    # Create a mock StepInput for step_generate_images
                    mock_input = StepInput(
                        input=user_prompt,
                        previous_step_content=plan_json,
                    )
                    step_generate_images(mock_input, chunk_session)  # noqa: F405
                    print(
                        "[PROCESS] Chunk %d: images generated. Count: %d"
                        % (chunk_idx, len(chunk_session.get("generated_images", {})))
                    )
                else:
                    print("[PROCESS] Chunk %d: image planner returned no plan." % chunk_idx)

            except Exception as e:
                print(
                    "[PROCESS] Chunk %d: image pipeline failed: %s" % (chunk_idx, e)
                )
                if session_state.get("verbose"):
                    traceback.print_exc()

        # --- Template assembly ---
        if template_path and os.path.isfile(template_path):
            print("[PROCESS] Chunk %d: running template assembly..." % chunk_idx)
            try:
                # Propagate generated images back to chunk_session
                mock_assemble_input = StepInput(
                    input=session_state.get("user_prompt", ""),
                    previous_step_content=json.dumps(slides_data),
                )
                step_assemble_template(mock_assemble_input, chunk_session)  # noqa: F405

                assembled_output = chunk_session.get("output_path", assembled_path)
                if assembled_output and os.path.isfile(assembled_output):
                    current_file = assembled_output
                    print(
                        "[PROCESS] Chunk %d: assembled -> %s" % (chunk_idx, current_file)
                    )
                else:
                    print(
                        "[PROCESS] Chunk %d: template assembly produced no file; "
                        "keeping raw chunk." % chunk_idx
                    )
            except Exception as e:
                print(
                    "[PROCESS] Chunk %d: template assembly failed: %s" % (chunk_idx, e)
                )
                if session_state.get("verbose"):
                    traceback.print_exc()
        else:
            # No template: just copy raw chunk to assembled path name for consistency
            shutil.copy2(chunk_file, assembled_path)
            current_file = assembled_path
            print(
                "[PROCESS] Chunk %d: no template; raw chunk copied to %s"
                % (chunk_idx, assembled_path)
            )

        processed_chunks[chunk_idx] = current_file

        chunk_proc_elapsed = time.time() - chunk_proc_start
        print("[TIMING] Chunk %d processing done in %.1fs" % (chunk_idx, chunk_proc_elapsed))
        print("[PROCESS] Chunk %d: result -> %s" % (chunk_idx, current_file))

    session_state["processed_chunks"] = processed_chunks
    successful = sum(1 for v in processed_chunks.values() if v is not None)

    step_elapsed = time.time() - step_start
    print(
        "\n[TIMING] step_process_chunks completed in %.1fs (%d chunks processed)"
        % (step_elapsed, successful)
    )

    return StepOutput(
        content="%d of %d chunks processed. Duration: %.1fs"
        % (successful, len(chunk_files), step_elapsed),
        success=True,
    )


# === WORKFLOW STEP 4 (OPTIONAL): VISUAL REVIEW PER CHUNK ===


def step_visual_review_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 4 (Optional): Run visual inspection on each chunk's assembled PPTX.

    For each chunk:
    1. Render slides to PNG via LibreOffice.
    2. Call slide_quality_reviewer for each slide image.
    3. Apply programmatic corrections if needed.
    4. Repeat up to max_passes passes until no further changes are needed.

    This step is non-blocking: any failure silently returns success=True.
    If a programmatic fix is missing in Python, logs it to console regardless
    of --verbose setting.
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 4 (Optional): Visual review per chunk...")
    print("=" * 60)

    processed_chunks: Dict[int, Optional[str]] = session_state.get(
        "processed_chunks", {}
    )
    output_dir = session_state.get("output_dir", ".")
    template_path = session_state.get("template_path", "")
    max_passes = session_state.get("visual_passes", 3)
    reviewed_chunks: Dict[int, Optional[str]] = {}

    for chunk_idx, assembled_path in sorted(processed_chunks.items()):
        chunk_review_start = time.time()

        if assembled_path is None or not os.path.isfile(assembled_path):
            print(
                "[VISUAL] Chunk %d: skipped (file not found: %s)."
                % (chunk_idx, assembled_path)
            )
            reviewed_chunks[chunk_idx] = None
            continue

        print(
            "\n[VISUAL REVIEW] Chunk %d: starting review of %s"
            % (chunk_idx, assembled_path)
        )

        # Build a per-chunk session_state for the visual review step
        chunk_session = dict(session_state)
        chunk_session["output_path"] = assembled_path
        chunk_session["template_path"] = template_path

        current_path = assembled_path

        for pass_num in range(max_passes):
            pass_start = time.time()
            print(
                "[VISUAL REVIEW] Chunk %d: pass %d/%d starting..."
                % (chunk_idx, pass_num + 1, max_passes)
            )
            chunk_session["output_path"] = current_path

            try:
                mock_input = StepInput(
                    input=session_state.get("user_prompt", ""),
                    previous_step_content="",
                )
                result = step_visual_quality_review(mock_input, chunk_session)  # noqa: F405

                # Check if any actionable issues exist (mirrors _apply_visual_corrections logic).
                # PresentationQualityReport has 'total_critical_issues', NOT 'total_corrections_applied'.
                # We check slide_reports for critical/moderate issues with a real programmatic_fix.
                quality_report = chunk_session.get("quality_report", {})
                slide_reports_data = quality_report.get("slide_reports", [])

                if VERBOSE:  # noqa: F405
                    for r in slide_reports_data:
                        issues = r.get("issues", [])
                        print(
                            "[VERBOSE] Chunk %d pass %d slide %s: %d issues"
                            % (
                                chunk_idx,
                                pass_num + 1,
                                r.get("slide_index", "?"),
                                len(issues),
                            )
                        )
                        for issue in issues:
                            print(
                                "[VERBOSE]   severity=%s fix=%s desc=%s"
                                % (
                                    issue.get("severity", "?"),
                                    issue.get("programmatic_fix", "?"),
                                    str(issue.get("description", ""))[:80],
                                )
                            )

                changes_applied = any(
                    any(
                        i.get("severity") in ("critical", "moderate")
                        and i.get("programmatic_fix") != "none"
                        for i in r.get("issues", [])
                    )
                    for r in slide_reports_data
                )

                pass_elapsed = time.time() - pass_start
                print(
                    "[TIMING] Chunk %d pass %d: %.1fs"
                    % (chunk_idx, pass_num + 1, pass_elapsed)
                )

                if not changes_applied:
                    print(
                        "[VISUAL REVIEW] Chunk %d: pass %d/%d — no changes needed. Done."
                        % (chunk_idx, pass_num + 1, max_passes)
                    )
                    break
                else:
                    print(
                        "[VISUAL REVIEW] Chunk %d: pass %d/%d — corrections applied. Re-checking..."
                        % (chunk_idx, pass_num + 1, max_passes)
                    )

            except Exception as e:
                pass_elapsed = time.time() - pass_start
                print(
                    "[TIMING] Chunk %d pass %d: %.1fs (error)"
                    % (chunk_idx, pass_num + 1, pass_elapsed)
                )
                print(
                    "[VISUAL] Chunk %d, pass %d: review failed: %s"
                    % (chunk_idx, pass_num + 1, e)
                )
                # Log missing programmatic fix to console regardless of verbose mode (per spec).
                # Any exception here means the visual review or correction logic is broken/missing.
                print(
                    "[VISUAL REVIEW MISSING FIX] Chunk %d, pass %d: exception during "
                    "visual correction step: %s"
                    % (chunk_idx, pass_num + 1, str(e))
                )
                print(
                    "[SUGGESTION] Review step_visual_quality_review() and "
                    "_apply_visual_corrections() for the issue type that raised this error. "
                    "Add handling logic if a programmatic_fix type is missing."
                )
                break

        reviewed_chunks[chunk_idx] = chunk_session.get("output_path", current_path)

        chunk_review_elapsed = time.time() - chunk_review_start
        print(
            "[TIMING] Chunk %d total review: %.1fs" % (chunk_idx, chunk_review_elapsed)
        )
        print(
            "[VISUAL REVIEW] Chunk %d: reviewed -> %s"
            % (chunk_idx, reviewed_chunks[chunk_idx])
        )

    session_state["reviewed_chunks"] = reviewed_chunks
    reviewed_count = sum(1 for v in reviewed_chunks.values() if v is not None)

    step_elapsed = time.time() - step_start
    print(
        "\n[TIMING] step_visual_review_chunks completed in %.1fs (%d chunks reviewed)"
        % (step_elapsed, reviewed_count)
    )

    return StepOutput(
        content="%d of %d chunks visually reviewed. Duration: %.1fs"
        % (reviewed_count, len(processed_chunks), step_elapsed),
        success=True,
    )


# === HELPER: MERGE MULTIPLE PPTX FILES ===


def _merge_pptx_zip_level(pptx_paths: List[str], output_path: str) -> bool:
    """Merge multiple PPTX files by manipulating the ZIP structure directly.

    This is the most reliable approach — avoids OPC context issues that cause
    PowerPoint to report "found a problem with content" on the merged file.
    Binary parts (images, charts, workbooks) are copied at the raw bytes level,
    so there is no risk of the OPC package context dropping or corrupting data.

    Args:
        pptx_paths: List of valid, existing PPTX file paths to merge in order.
        output_path: Destination path for the merged presentation.

    Returns:
        True if merge succeeded, False otherwise.
    """
    import re
    from lxml import etree

    valid_paths = [p for p in pptx_paths if p and os.path.exists(p)]
    if not valid_paths:
        print("[MERGE] No valid PPTX files to merge")
        return False

    if len(valid_paths) == 1:
        shutil.copy(valid_paths[0], output_path)
        print("[MERGE] Single file, copied directly: %s" % output_path)
        return True

    # Use first file as base — copy it to output
    shutil.copy(valid_paths[0], output_path)

    # Read the base presentation XML
    with zipfile.ZipFile(output_path, "r") as base_zip:
        base_prs_xml = base_zip.read("ppt/presentation.xml")
        base_prs_rels_xml = base_zip.read("ppt/_rels/presentation.xml.rels")
        base_content_types_xml = base_zip.read("[Content_Types].xml")

    base_prs_tree = etree.fromstring(base_prs_xml)
    base_prs_rels_tree = etree.fromstring(base_prs_rels_xml)
    base_ct_tree = etree.fromstring(base_content_types_xml)

    # Namespaces
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_RELS = "http://schemas.openxmlformats.org/package/2006/relationships"
    NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"

    def _get_slide_numbers(prs_rels_tree):
        slide_nums = []
        for rel in prs_rels_tree.findall("{%s}Relationship" % NS_RELS):
            target = rel.get("Target", "")
            m = re.match(r"slides/slide(\d+)\.xml", target)
            if m:
                slide_nums.append(int(m.group(1)))
        return slide_nums

    existing_slide_nums = _get_slide_numbers(base_prs_rels_tree)
    next_slide_num = max(existing_slide_nums) + 1 if existing_slide_nums else 1

    # Track next rel ID for presentation.xml.rels
    existing_rel_ids = [
        int(rel.get("Id", "rId0").replace("rId", ""))
        for rel in base_prs_rels_tree.findall("{%s}Relationship" % NS_RELS)
        if rel.get("Id", "").startswith("rId")
    ]
    next_rel_id = max(existing_rel_ids) + 1 if existing_rel_ids else 100

    # Open output as writable archive
    with zipfile.ZipFile(output_path, "a", compression=zipfile.ZIP_DEFLATED, allowZip64=True) as out_zip:
        for src_path in valid_paths[1:]:
            with zipfile.ZipFile(src_path, "r") as src_zip:
                src_names = set(src_zip.namelist())

                # Read source presentation rels to find slides
                src_prs_rels_xml = src_zip.read("ppt/_rels/presentation.xml.rels")
                src_prs_rels = etree.fromstring(src_prs_rels_xml)
                src_slide_nums = sorted(_get_slide_numbers(src_prs_rels))

                for src_slide_num in src_slide_nums:
                    new_slide_num = next_slide_num
                    next_slide_num += 1

                    old_slide_name = "ppt/slides/slide%d.xml" % src_slide_num
                    new_slide_name = "ppt/slides/slide%d.xml" % new_slide_num
                    old_slide_rels_name = "ppt/slides/_rels/slide%d.xml.rels" % src_slide_num
                    new_slide_rels_name = "ppt/slides/_rels/slide%d.xml.rels" % new_slide_num

                    if old_slide_name not in src_names:
                        continue

                    # Copy slide XML
                    slide_xml_bytes = src_zip.read(old_slide_name)
                    out_zip.writestr(new_slide_name, slide_xml_bytes)

                    # Copy slide rels and rewrite media/chart refs with unique names
                    if old_slide_rels_name in src_names:
                        slide_rels_xml = src_zip.read(old_slide_rels_name)
                        slide_rels_tree = etree.fromstring(slide_rels_xml)

                        for rel in slide_rels_tree.findall("{%s}Relationship" % NS_RELS):
                            rel_type = rel.get("Type", "")
                            target = rel.get("Target", "")

                            # Skip slide layout — keep reference as-is
                            if "slideLayout" in rel_type:
                                continue

                            if not target.startswith(".."):
                                continue  # absolute or external refs

                            # Resolve the actual part path in the source zip
                            # e.g. "../media/image1.png" -> "ppt/media/image1.png"
                            actual_old = "ppt/slides/" + target.lstrip("../")

                            if actual_old not in src_names:
                                continue

                            # Generate unique name for target archive
                            basename = os.path.basename(actual_old)
                            stem, ext = os.path.splitext(basename)
                            new_part_name = actual_old
                            counter = 1
                            all_names = set(out_zip.namelist())
                            while new_part_name in all_names:
                                new_part_name = (
                                    os.path.dirname(actual_old)
                                    + "/"
                                    + stem
                                    + "_s%d_%d" % (new_slide_num, counter)
                                    + ext
                                )
                                counter += 1

                            # Copy the part
                            part_bytes = src_zip.read(actual_old)
                            out_zip.writestr(new_part_name, part_bytes)

                            # If chart, also copy its rels and embedded workbook
                            if "chart" in actual_old:
                                chart_basename = os.path.basename(actual_old)
                                chart_rels_old = (
                                    os.path.dirname(actual_old)
                                    + "/_rels/"
                                    + chart_basename
                                    + ".rels"
                                )
                                if chart_rels_old in src_names:
                                    cr_bytes = src_zip.read(chart_rels_old)
                                    cr_new = (
                                        os.path.dirname(new_part_name)
                                        + "/_rels/"
                                        + os.path.basename(new_part_name)
                                        + ".rels"
                                    )
                                    out_zip.writestr(cr_new, cr_bytes)
                                    # Copy chart's embedded xlsx workbook
                                    cr_tree = etree.fromstring(cr_bytes)
                                    for cr_rel in cr_tree.findall("{%s}Relationship" % NS_RELS):
                                        cr_target = cr_rel.get("Target", "")
                                        if cr_target.startswith(".."):
                                            wb_old = "ppt/charts/" + cr_target.lstrip("../")
                                            if wb_old in src_names:
                                                wb_bytes = src_zip.read(wb_old)
                                                wb_stem, wb_ext = os.path.splitext(
                                                    os.path.basename(wb_old)
                                                )
                                                wb_new = wb_old.replace(
                                                    os.path.basename(wb_old),
                                                    wb_stem
                                                    + "_s%d_wb" % new_slide_num
                                                    + wb_ext,
                                                )
                                                if wb_new not in set(out_zip.namelist()):
                                                    out_zip.writestr(wb_new, wb_bytes)

                            # Update the relationship target to point to new unique name
                            new_rel_target = "../" + "/".join(new_part_name.split("/")[2:])
                            rel.set("Target", new_rel_target)

                        updated_rels_bytes = etree.tostring(
                            slide_rels_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        )
                        out_zip.writestr(new_slide_rels_name, updated_rels_bytes)

                    # Register slide in presentation.xml.rels
                    new_rel_id = "rId%d" % next_rel_id
                    next_rel_id += 1
                    new_prs_rel = etree.SubElement(
                        base_prs_rels_tree, "{%s}Relationship" % NS_RELS
                    )
                    new_prs_rel.set("Id", new_rel_id)
                    new_prs_rel.set(
                        "Type",
                        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
                    )
                    new_prs_rel.set("Target", "slides/slide%d.xml" % new_slide_num)

                    # Register slide in presentation.xml sldIdLst
                    sld_id_lst = base_prs_tree.find(".//{%s}sldIdLst" % NS_P)
                    if sld_id_lst is None:
                        sld_id_lst = etree.SubElement(
                            base_prs_tree, "{%s}sldIdLst" % NS_P
                        )

                    existing_ids = [
                        int(el.get("id", 256))
                        for el in sld_id_lst.findall("{%s}sldId" % NS_P)
                    ]
                    new_id = max(existing_ids) + 1 if existing_ids else 256
                    sld_id_el = etree.SubElement(sld_id_lst, "{%s}sldId" % NS_P)
                    sld_id_el.set("id", str(new_id))
                    sld_id_el.set("{%s}id" % NS_R, new_rel_id)

                    # Add content type entry for new slide
                    existing_ct_parts = {
                        el.get("PartName", "")
                        for el in base_ct_tree.findall("{%s}Override" % NS_CT)
                    }
                    new_part_uri = "/ppt/slides/slide%d.xml" % new_slide_num
                    if new_part_uri not in existing_ct_parts:
                        ct_el = etree.SubElement(base_ct_tree, "{%s}Override" % NS_CT)
                        ct_el.set("PartName", new_part_uri)
                        ct_el.set(
                            "ContentType",
                            "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
                        )

    # Write updated presentation.xml, rels, and content types back.
    # Python's zipfile does not support overwriting entries, so we copy to a temp file.
    import tempfile
    tmp_path = output_path + ".tmp"
    with zipfile.ZipFile(output_path, "r") as old_zip:
        with zipfile.ZipFile(
            tmp_path, "w", compression=zipfile.ZIP_DEFLATED, allowZip64=True
        ) as new_zip:
            for item in old_zip.namelist():
                if item == "ppt/presentation.xml":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_prs_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                elif item == "ppt/_rels/presentation.xml.rels":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_prs_rels_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                elif item == "[Content_Types].xml":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_ct_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                else:
                    new_zip.writestr(item, old_zip.read(item))

    os.replace(tmp_path, output_path)
    print("[MERGE] Saved merged presentation: %s" % output_path)
    return True


def _try_auto_repair_with_libreoffice(pptx_path: str) -> bool:
    """Attempt to auto-repair a PPTX by converting it through LibreOffice.

    Only runs if LibreOffice is available. Non-destructive on failure —
    the original file is left unchanged if the conversion fails.

    Args:
        pptx_path: Path to the PPTX file to repair in-place.

    Returns:
        True if LibreOffice repair succeeded, False otherwise.
    """
    import glob
    import subprocess
    import tempfile

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return False

    tmp_dir = tempfile.mkdtemp(prefix="pptx_repair_")
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pptx", "--outdir", tmp_dir, pptx_path],
            capture_output=True,
            timeout=180,
        )
        if result.returncode == 0:
            converted = glob.glob(os.path.join(tmp_dir, "*.pptx"))
            if converted:
                shutil.copy(converted[0], pptx_path)
                print("[MERGE] Auto-repair via LibreOffice succeeded: %s" % pptx_path)
                return True
    except Exception as e:
        print("[MERGE] Auto-repair via LibreOffice failed: %s" % e)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
    return False


def merge_pptx_files(pptx_paths: List[str], output_path: str) -> bool:
    """Merge multiple PPTX files into a single presentation using ZIP-level manipulation.

    Uses _merge_pptx_zip_level() which copies all binary parts (images, charts,
    workbooks) at the raw bytes level, avoiding OPC package context issues that
    cause PowerPoint to report "found a problem with content" on the merged file.

    Args:
        pptx_paths: List of PPTX file paths to merge in order.
        output_path: Destination path for the merged presentation.

    Returns:
        True if merge succeeded, False otherwise.
    """
    merge_start = time.time()
    valid_paths = [p for p in pptx_paths if p and os.path.exists(p)]
    print("[MERGE] Merging %d PPTX files into %s" % (len(valid_paths), output_path))
    if VERBOSE:  # noqa: F405
        for i, p in enumerate(valid_paths):
            print("[VERBOSE][MERGE] Source %d: %s" % (i, p))
    result = _merge_pptx_zip_level(valid_paths, output_path)
    merge_elapsed = time.time() - merge_start
    print("[TIMING] merge_pptx_files completed in %.1fs" % merge_elapsed)
    return result


# === WORKFLOW STEP 5 (FINAL): MERGE ALL CHUNKS ===


def step_merge_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 5 (Final): Merge all processed/reviewed chunk PPTX files into the final output.

    Source selection priority (explicit, robust):
    1. Template + visual review + reviewed_chunks present -> use reviewed_chunks
    2. Template + processed_chunks present             -> use processed_chunks
    3. No template (raw mode)                          -> use raw chunk_files

    Chunks are merged in order (by chunk_idx).
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 5 (Final): Merging chunks into final presentation...")
    print("=" * 60)

    output_path = session_state.get("output_path", "presentation_chunked.pptx")
    has_template = bool(session_state.get("template_path"))
    visual_review = session_state.get("visual_review", False)
    chunk_files: List[Optional[str]] = session_state.get("chunk_files", [])
    processed_chunks: Dict[int, Optional[str]] = session_state.get("processed_chunks", {})
    reviewed_chunks: Dict[int, Optional[str]] = session_state.get("reviewed_chunks", {})

    # Determine which chunk paths to use (priority: reviewed > processed > raw)
    if has_template and visual_review and reviewed_chunks:
        source_label = "reviewed (template + visual review)"
        ordered_paths = [
            reviewed_chunks.get(i)
            for i in sorted(reviewed_chunks.keys())
        ]
    elif has_template and processed_chunks:
        source_label = "processed (template-assembled)"
        ordered_paths = [
            processed_chunks.get(i)
            for i in sorted(processed_chunks.keys())
        ]
    else:
        # No template path: use raw chunk files directly
        source_label = "raw (no template)"
        ordered_paths = [f for f in chunk_files if f is not None]

    if not ordered_paths:
        print("[MERGE] No chunk files found to merge")
        return StepOutput(
            content="No files to merge",
            success=False,
        )

    print(
        "Merging from: %s (%d total, %d valid)"
        % (
            source_label,
            len(ordered_paths),
            sum(1 for p in ordered_paths if p and os.path.exists(p)),
        )
    )

    if VERBOSE:  # noqa: F405
        print("[VERBOSE] Ordered chunk files for merge:")
        for i, p in enumerate(ordered_paths):
            print("[VERBOSE]   %d. %s" % (i, p))

    success = merge_pptx_files(
        [p for p in ordered_paths if p],
        output_path,
    )

    # Attempt optional auto-repair (only if LibreOffice is available)
    if success:
        _try_auto_repair_with_libreoffice(output_path)

    step_elapsed = time.time() - step_start
    final_file = os.path.basename(output_path)
    print(
        "[TIMING] step_merge_chunks completed in %.1fs (final: %s)"
        % (step_elapsed, final_file)
    )

    if success:
        summary = "Merged %d chunks (%s) -> %s. Duration: %.1fs" % (
            len([p for p in ordered_paths if p]),
            source_label,
            output_path,
            step_elapsed,
        )
        print("[MERGE] %s" % summary)
        return StepOutput(
            content=summary,
            success=True,
        )
    else:
        return StepOutput(
            content="Merge failed. No output file produced.",
            success=False,
        )


# === WORKFLOW BUILDER ===


def build_chunked_workflow(session_state: Dict) -> Workflow:
    """Build the chunked PPTX workflow with the appropriate set of steps.

    Steps included:
    - Step 1: Optimize & Plan   (always)
    - Step 2: Generate Chunks   (always)
    - Step 3: Process Chunks    (only when template_path is set)
    - Step 4: Visual Review     (only when template_path AND visual_review are both set)
    - Step 5: Merge Chunks      (always)

    No-template pipeline: Step 1 -> Step 2 -> Step 5
    Template pipeline:    Step 1 -> Step 2 -> Step 3 [-> Step 4] -> Step 5
    """
    has_template = bool(session_state.get("template_path"))
    do_visual_review = has_template and bool(session_state.get("visual_review"))

    steps = [
        Step(name="Optimize and Plan", executor=step_optimize_and_plan),
        Step(name="Generate Chunks", executor=step_generate_chunks),
    ]

    # Template assembly + image pipeline only runs when a template is provided
    if has_template:
        steps.append(Step(name="Process Chunks", executor=step_process_chunks))

    # Visual review requires both --visual-review flag AND a template
    if do_visual_review:
        steps.append(
            Step(name="Visual Review Chunks", executor=step_visual_review_chunks)
        )

    steps.append(Step(name="Merge Chunks", executor=step_merge_chunks))

    return Workflow(
        name="Chunked PPTX Workflow",
        steps=steps,
        session_state=session_state,
    )


# === MAIN ENTRY POINT ===


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Chunked PPTX generation workflow — overcomes Claude API limits for large presentations."
    )

    # Existing args (compatible with powerpoint_template_workflow.py)
    parser.add_argument(
        "--template",
        "-t",
        default=None,
        help="Path to .pptx template file (optional). Without it, skips template assembly.",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="presentation_chunked.pptx",
        help="Output filename (default: presentation_chunked.pptx).",
    )
    parser.add_argument(
        "--prompt",
        "-p",
        default=None,
        help="User prompt describing the presentation topic.",
    )
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Skip AI image generation.",
    )
    parser.add_argument(
        "--no-stream",
        action="store_true",
        help="Disable streaming mode for Claude agent.",
    )
    parser.add_argument(
        "--min-images",
        type=int,
        default=1,
        help="Minimum slides that must have AI-generated images (default: 1).",
    )
    parser.add_argument(
        "--visual-review",
        action="store_true",
        help="Enable visual QA with Gemini vision per chunk (requires LibreOffice + template).",
    )
    parser.add_argument(
        "--footer-text",
        default="",
        help="Footer text for all slides (idx=11 placeholder).",
    )
    parser.add_argument(
        "--date-text",
        default="",
        help="Date text for footer date placeholder (idx=10).",
    )
    parser.add_argument(
        "--show-slide-numbers",
        action="store_true",
        help="Preserve slide number placeholder (idx=12) on all slides.",
    )
    parser.add_argument(
        "--verbose",
        "-v",
        action="store_true",
        help="Enable verbose/debug logging.",
    )

    # New args for chunked workflow
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=3,
        help="Number of slides per Claude API chunk call (default: 3).",
    )
    parser.add_argument(
        "--max-retries",
        type=int,
        default=2,
        help="Max retries per chunk on failure (default: 2).",
    )
    parser.add_argument(
        "--visual-passes",
        type=int,
        default=3,
        help="Maximum visual inspection passes per chunk (default: 3).",
    )

    args = parser.parse_args()

    # Update module-level VERBOSE (imported from powerpoint_template_workflow via *)
    global VERBOSE  # noqa: F405
    VERBOSE = args.verbose  # noqa: F405

    # Validate API key
    if not os.getenv("ANTHROPIC_API_KEY"):
        print("Error: ANTHROPIC_API_KEY environment variable not set.")
        sys.exit(1)

    # Validate template if provided
    if args.template is not None:
        if not os.path.isfile(args.template):
            print("Error: Template file not found: %s" % args.template)
            sys.exit(1)
        if not args.template.endswith(".pptx"):
            print("Error: Template file must be a .pptx file.")
            sys.exit(1)

    # Warn when visual flags are passed without a template (they will be ignored)
    if not args.template and args.visual_review:
        print("[WARNING] --visual-review is ignored when --template is not provided")
    if not args.template and args.visual_passes != 3:
        print("[WARNING] --visual-passes is ignored when --template is not provided")

    # Effective values: visual review and passes are forced off when no template
    effective_visual_review = bool(args.visual_review) and bool(args.template)
    effective_visual_passes = args.visual_passes if args.template else 0

    # Setup output directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_base = os.path.join(script_dir, "output_chunked")
    os.makedirs(output_base, exist_ok=True)

    # Resolve output path
    output_path = args.output
    if not os.path.isabs(output_path):
        output_path = os.path.join(output_base, output_path)

    # Chunked workflow uses its own working directory
    output_dir = os.path.join(output_base, "chunked_workflow_work")
    os.makedirs(output_dir, exist_ok=True)

    default_prompt = (
        "Create a professional business presentation about AI transformation in enterprise companies"
    )

    session_state = {
        # Core paths
        "template_path": args.template or "",
        "output_path": output_path,
        "output_dir": output_dir,
        # User inputs
        "user_prompt": args.prompt or default_prompt,
        "verbose": args.verbose,
        "stream": not args.no_stream,
        "no_images": args.no_images,
        "min_images": args.min_images,
        # visual_review and visual_passes are forced False/0 when no template
        "visual_review": effective_visual_review,
        "footer_text": args.footer_text,
        "date_text": args.date_text,
        "show_slide_numbers": args.show_slide_numbers,
        # Chunked workflow settings
        "chunk_size": args.chunk_size,
        "max_retries": args.max_retries,
        "visual_passes": effective_visual_passes,
        # Fields populated by steps
        "storyboard": None,
        "storyboard_dir": None,
        "total_slides": 0,
        "chunk_files": [],
        "chunk_slide_groups": [],
        "processed_chunks": {},
        "reviewed_chunks": {},
        # Fields used by existing step helpers
        "generated_file": "",
        "slides_data": [],
        "generated_images": {},
        "src_slide_width": 0,
        "src_slide_height": 0,
        "assembly_knowledge": {},
        "quality_report": {},
    }

    workflow = build_chunked_workflow(session_state)

    print("=" * 60)
    print("Chunked PPTX Workflow")
    print("=" * 60)
    print("Prompt:     %s" % (args.prompt or default_prompt)[:80])
    print("Output:     %s" % output_path)
    if args.template:
        print("Mode:       template-assisted generation")
        print("Template:   %s" % args.template)
        if effective_visual_review:
            print("Visual review: enabled (%d passes max)" % args.visual_passes)
        else:
            print("Visual review: disabled")
    else:
        print("Mode:       raw generation (no template)")
        print("Visual review: skipped (no template)")
    print("Chunk size: %d slides per API call" % args.chunk_size)
    print("Max retries per chunk: %d" % args.max_retries)
    print("Images:     %s" % ("disabled" if args.no_images else "enabled"))
    if args.verbose:
        print("Verbose:    enabled")

    start_time = time.time()

    workflow.run()

    elapsed = time.time() - start_time
    print("\n" + "=" * 60)
    print("[TIMING] Total workflow: %.1fs" % elapsed)
    print("Output: %s" % output_path)
    print("=" * 60)


if __name__ == "__main__":
    main()
