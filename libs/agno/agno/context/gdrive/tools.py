"""
Drive toolkit that actually sees shared folders and Shared Drives.

The upstream ``GoogleDriveTools.search_files`` / ``_get_file_metadata`` /
``read_file`` methods don't pass ``corpora="allDrives"``,
``includeItemsFromAllDrives=True``, or ``supportsAllDrives=True``. That
defaults the Drive API to ``corpora=user``, which for a service account
means "files directly owned by the SA" — nothing shared with it and
nothing in Shared Drives. So a standard service-account setup (SA +
folders shared by humans + files in Shared Drives) returns zero hits
for the most ordinary ``name contains 'X'`` query.

This subclass overrides the three call sites and injects the allDrives
triple on every request. Everything else (auth, ``include_trashed``,
field selection, error handling) is inherited unchanged.

Kept local to ``agno.context.gdrive`` instead of fixing upstream so
callers of ``GoogleDriveTools`` directly aren't affected by the scope
change.
"""

from __future__ import annotations

import io
import json
from typing import Any, Optional, cast

from agno.tools.google.drive import (
    GoogleDriveTools,
    WorkspaceType,
    _is_binary_mime,
    authenticate,
)
from agno.utils.log import log_error

try:
    from googleapiclient.discovery import Resource
    from googleapiclient.errors import HttpError
    from googleapiclient.http import MediaIoBaseDownload
except ImportError:
    raise ImportError(
        "Google client library for Python not found, install it using "
        "`pip install google-api-python-client google-auth-httplib2 google-auth-oauthlib`"
    )


# Office formats we can extract text from with optional dependencies
DOCX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


# Optional dependency extraction helpers
# Each function imports its library inline and raises ImportError if missing.
# The caller (read_file) catches ImportError and returns a user-friendly install message.
# This keeps the helpers reusable and error handling centralized.


def _extract_docx_text(content_bytes: bytes) -> str:
    """Extract text content from a .docx file. Raises ImportError if python-docx not installed."""
    import docx  # inline import — caller handles ImportError

    buffer = io.BytesIO(content_bytes)
    document = docx.Document(buffer)
    paragraphs = [p.text for p in document.paragraphs]
    return "\n".join(paragraphs)


def _extract_xlsx_text(content_bytes: bytes) -> str:
    """Extract text content from a .xlsx file. Raises ImportError if openpyxl not installed."""
    import openpyxl  # inline import — caller handles ImportError

    buffer = io.BytesIO(content_bytes)
    workbook = openpyxl.load_workbook(buffer, read_only=True, data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        lines.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    return "\n".join(lines)


def _extract_pptx_text(content_bytes: bytes) -> str:
    """Extract text content from a .pptx file. Raises ImportError if python-pptx not installed."""
    from pptx import Presentation  # type: ignore[import-not-found]

    buffer = io.BytesIO(content_bytes)
    prs = Presentation(buffer)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        lines.append(text)
    return "\n".join(lines)


class AllDrivesGoogleDriveTools(GoogleDriveTools):
    """Drive toolkit that searches personal + shared + Shared Drive corpora."""

    @authenticate
    def search_files(self, query: Optional[str] = None, max_results: int = 10, page_token: Optional[str] = None) -> str:
        if max_results < 1:
            return json.dumps({"error": "max_results must be greater than 0"})
        try:
            service = cast(Resource, self.service)
            if self.include_trashed:
                effective_query = query or ""
            elif query:
                effective_query = f"({query}) and trashed=false"
            else:
                effective_query = "trashed=false"
            list_kwargs: dict = {
                "q": effective_query,
                "pageSize": max_results,
                "orderBy": "modifiedTime desc",
                "fields": self.SEARCH_FIELDS,
                "corpora": "allDrives",
                "includeItemsFromAllDrives": True,
                "supportsAllDrives": True,
            }
            if page_token:
                list_kwargs["pageToken"] = page_token
            results = service.files().list(**list_kwargs).execute()
            files = results.get("files", [])
            return json.dumps(
                {
                    "query": effective_query,
                    "files": files,
                    "count": len(files),
                    "nextPageToken": results.get("nextPageToken"),
                }
            )
        except HttpError as e:
            return json.dumps({"error": f"Google Drive API error: {e}"})
        except Exception as e:
            log_error(f"Could not search Google Drive files: {e}")
            return json.dumps({"error": f"Unexpected error: {type(e).__name__}: {e}"})

    def _get_file_metadata(self, file_id: str, fields: str) -> dict:
        service = cast(Resource, self.service)
        return service.files().get(fileId=file_id, fields=fields, supportsAllDrives=True).execute()

    @authenticate
    def read_file(self, file_id: str) -> str:
        try:
            service = cast(Resource, self.service)
            metadata = self._get_file_metadata(file_id, self.READ_METADATA_FIELDS)
            mime_type = metadata.get("mimeType", "")

            if mime_type in self.TEXT_EXPORT_TYPES:
                export_mime: Optional[str] = self.TEXT_EXPORT_TYPES[mime_type]
            elif mime_type.startswith(WorkspaceType.WORKSPACE_PREFIX):
                return json.dumps(
                    {"error": f"Cannot read {mime_type} as text. Use download_file instead.", "file": metadata}
                )
            elif mime_type in (DOCX_MIME_TYPE, XLSX_MIME_TYPE, PPTX_MIME_TYPE):
                file_size = int(metadata.get("size", 0))
                if file_size > self.max_read_size:
                    return json.dumps(
                        {
                            "error": (
                                f"File is {file_size} bytes, exceeds max_read_size "
                                f"({self.max_read_size}). Use download_file instead."
                            ),
                            "file": metadata,
                        }
                    )
                request = service.files().get_media(fileId=file_id, supportsAllDrives=True)
                content_bytes = _download_bytes(request)
                fmt, pkg = {
                    DOCX_MIME_TYPE: ("docx", "python-docx"),
                    XLSX_MIME_TYPE: ("xlsx", "openpyxl"),
                    PPTX_MIME_TYPE: ("pptx", "python-pptx"),
                }[mime_type]
                try:
                    if mime_type == DOCX_MIME_TYPE:
                        content = _extract_docx_text(content_bytes)
                    elif mime_type == XLSX_MIME_TYPE:
                        content = _extract_xlsx_text(content_bytes)
                    else:
                        content = _extract_pptx_text(content_bytes)
                    return json.dumps(
                        {
                            "file": metadata,
                            "content": content,
                            "contentLength": len(content),
                            "extractedFrom": fmt,
                        }
                    )
                except ImportError:
                    return json.dumps(
                        {
                            "error": f"Cannot read .{fmt} file: {pkg} not installed. Install with: pip install {pkg}",
                            "file": metadata,
                        }
                    )
            elif _is_binary_mime(mime_type):
                file_ext = metadata.get("name", "").rsplit(".", 1)[-1].lower()
                hint = ""
                if file_ext in ("docx", "doc"):
                    hint = " To read as text, open in Google Drive and convert to Google Docs format."
                elif file_ext in ("xlsx", "xls"):
                    hint = " To read as text, open in Google Drive and convert to Google Sheets format."
                elif file_ext in ("pptx", "ppt"):
                    hint = " To read as text, open in Google Drive and convert to Google Slides format."
                return json.dumps(
                    {
                        "error": f"Cannot read binary file ({mime_type}) as text.{hint}",
                        "file": metadata,
                    }
                )
            else:
                export_mime = None

            if export_mime:
                request = service.files().export_media(fileId=file_id, mimeType=export_mime)
                content_bytes = _download_bytes(request)
            else:
                file_size = int(metadata.get("size", 0))
                if file_size > self.max_read_size:
                    return json.dumps(
                        {
                            "error": (
                                f"File is {file_size} bytes, exceeds max_read_size "
                                f"({self.max_read_size}). Use download_file instead."
                            ),
                            "file": metadata,
                        }
                    )
                request = service.files().get_media(fileId=file_id, supportsAllDrives=True)
                content_bytes = _download_bytes(request)

            content = content_bytes.decode("utf-8", errors="replace")
            return json.dumps(
                {
                    "file": metadata,
                    "content": content,
                    "contentLength": len(content),
                    "exportMimeType": export_mime,
                }
            )
        except HttpError as e:
            return json.dumps({"error": f"Google Drive API error: {e}"})
        except Exception as e:
            log_error(f"Could not read Google Drive file {file_id}: {e}")
            return json.dumps({"error": f"Unexpected error: {type(e).__name__}: {e}"})


def _download_bytes(request: Any) -> bytes:
    buffer = io.BytesIO()
    downloader = MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buffer.getvalue()
